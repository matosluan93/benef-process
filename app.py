import os
import uuid
import re
import json
import zipfile
import unicodedata
from datetime import date
from ppt_generator import generate as generate_ppt
from flask import Flask, request, jsonify, render_template, send_file
import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 500 * 1024 * 1024  # 500MB

TEMP_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'temp')
os.makedirs(TEMP_DIR, exist_ok=True)

# ─── Configuração de colunas ──────────────────────────────────────────────────
# ORDEM IMPORTA: auto_map_cols para no primeiro alias que bater.
# Aliases mais específicos (descritivos) vêm antes de aliases curtos/ambíguos.

COL_ALIASES = {
    'empresa': [
        'nome da empresa', 'nome empresa',
        'empresa', 'company', 'razao social', 'unidade',
    ],
    'cnpj': [
        'cnpj da empresa', 'cnpj empresa', 'cnpj filial',
        'cnpj', 'cnpj/cpf',
    ],
    'nome': [
        'nome profissional', 'nome do profissional',
        'nome', 'colaborador', 'funcionario',
    ],
    'matricula': [
        'matrícula', 'matricula',
        'matrícula do profissional', 'matricula do profissional',
        'chapa', 'registro',
    ],
    'email': [
        # 'email funcional' (sem hífen) cobre IHM cujo cabeçalho é 'EMAIL FUNCIONAL'
        'email funcional', 'e-mail funcional',
        'email', 'e-mail',
        'email profissional', 'e-mail profissional',
    ],
    'tipo_vinculo': [
        # FIX: Stefanini tem 'Desc. Tipo de Vínculo' = 'CONTRATADO'/'ESTAGIÁRIO'.
        # A coluna 'Vínculo' contém apenas o código de letra ('C'), que nunca
        # bateria nas keywords de ESTAGIO_KW / PJ_KW.
        'desc. tipo de vínculo', 'desc. tipo de vinculo',
        'desc. vínculo', 'desc. vinculo',
        'vínculo', 'vinculo',
        'tipo de vínculo', 'tipo de vinculo',
        'regime', 'modalidade',
    ],
    'data_admissao': [
        'admissão', 'admissao',
        'data de admissão', 'data de admissao',
        'dt admissão', 'dt. admissão', 'dt admissao', 'dt. admissao',
        'data admissão', 'data admissao',
        'data_admissao', 'data_admissão',
    ],
    'desconto_folha': [
        'desconto em folha', 'desconto folha',
        'aceite wellhub', 'wellhub folha',
    ],
    'cargo': [
        # FIX: 'nome cargo' e 'nome função' primeiro para evitar que 'função'
        # (alias curto) bata na coluna de CÓDIGO NUMÉRICO 'Função' da Stefanini
        # ('8449') em vez do campo descritivo 'Nome Cargo' / 'Nome Função'.
        'nome cargo', 'nome função', 'nome funcao',
        'função do profissional', 'funcao do profissional',
        'função', 'funcao', 'cargo',
    ],
    'cpf': [
        'cpf do profissional', 'número de cpf', 'numero de cpf',
        'cpf', 'nr cpf',
    ],
}

# ─── Regex ────────────────────────────────────────────────────────────────────

EMAIL_RE  = re.compile(r'^[^\s@]+@[^\s@]+\.[^\s@]+$')
ISO_DATE  = re.compile(r'^\d{4}[-/]\d{1,2}[-/]\d{1,2}')   # YYYY-MM-DD ou YYYY/MM/DD
CPF_CLEAN = re.compile(r'\D')

ALLOWED_EMAIL_DOMAINS = {
    'ar.stefanini.com',
    'bankinginabox.site',
    'caps.haus',
    'cl.stefanini.com',
    'co.stefanini.com',
    'datastorm.com.br',
    'ecglobal.com',
    'ecglobal.haus',
    'ecossistema.haus',
    'gauge.com.br',
    'gauge.haus',
    'huia.haus',
    'ihm.com.br',
    'ihmsenior.com.br',
    'inspiring.haus',
    'latam.stefanini.com',
    'marcostefanini.com',
    'marcostefanini.com.br',
    'mx.stefanini.com',
    'necxt.com.br',
    'necxtorbitall.com.br',
    'openstartups.com.br',
    'orbitall.com.br',
    'orbitallpay.com.br',
    'originacao.com',
    'originacao.com.br',
    'pe.stefanini.com',
    'perep.eu',
    'perepanalytics.eu',
    'pontocertificado.com',
    'pontocertificado.com.br',
    'potenciaisstefanini.com.br',
    'reseller.stefanini.com',
    'sbc091.teamsonevoice.com',
    'scalait.com',
    'seniorengenharia.com.br',
    'singulahr.com.br',
    'sophie.chat',
    'stefanini.com',
    'stefanini.com.br',
    'stefanini.org.br',
    'stefaniniathome.com.br',
    'stefaninicyber.com',
    'stefaninilatam.mail.onmicrosoft.com',
    'stefaninilatam.onmicrosoft.com',
    'stefaninirafael.com',
    'stefaninirafael.com.br',
    'stefaniniservico.com.br',
    'stefaninitrends.com',
    'stefaninitrends.com.br',
    'stfeacesso.com.br',
    'sunrising.com.br',
    'sv.stefanini.com',
    'techteam.biz',
    'techteam.com',
    'useniu.com.br',
    'w3.haus',
}

# ─── Funções de normalização ──────────────────────────────────────────────────

def norm(v):
    """Lowercase + strip + remove BOM invisível (\\ufeff) de cabeçalhos ERP."""
    return str(v or '').strip().lstrip('\ufeff').lower()

def norm_ascii(v):
    """Normaliza texto para comparações sem acento, mantendo símbolos comuns."""
    txt = norm(v).replace('ï»¿', '')
    return ''.join(
        ch for ch in unicodedata.normalize('NFKD', txt)
        if not unicodedata.combining(ch)
    )

def clean_val(v):
    """
    Retorna string limpa. Trata isna() com try/except para não quebrar
    em tipos ambíguos (arrays, objetos customizados) que surgem com na_filter=False.
    """
    try:
        if pd.isna(v):
            return ''
    except (TypeError, ValueError):
        pass
    v_str = str(v).replace('\xa0', ' ').strip()
    if v_str.lower() in ('nan', 'none', 'nat', 'null', '-', ''):
        return ''
    return v_str

def format_cpf(v):
    """
    FIX: CPF da Topaz chega sem formatação ('61213829020').
    Formata para '612.138.290-20' quando há exatamente 11 dígitos.
    """
    raw = CPF_CLEAN.sub('', str(v))
    if len(raw) == 11:
        return f'{raw[:3]}.{raw[3:6]}.{raw[6:9]}-{raw[9:]}'
    return clean_val(v)


def format_cnpj(v):
    """Formata CNPJ como XX.XXX.XXX/XXXX-XX, corrigindo zeros à esquerda."""
    digits = ''.join(filter(str.isdigit, str(v or '')))
    if not digits:
        return str(v or '').strip()
    digits = digits.zfill(14)
    if len(digits) == 14:
        return f"{digits[:2]}.{digits[2:5]}.{digits[5:8]}/{digits[8:12]}-{digits[12:]}"
    return str(v or '').strip()


# ─── Parsing de datas ─────────────────────────────────────────────────────────

def parse_mixed_dates(s):
    try:
        if pd.isna(s) or str(s).strip() == '': return pd.NaT
    except Exception: pass
    val = str(s).strip()
    if val.replace('.', '', 1).isdigit():
        try: return pd.to_datetime(float(val), origin='1899-12-30', unit='D')
        except Exception: pass
    if re.match(r'^\d{4}[-/]\d', val):
        return pd.to_datetime(val, dayfirst=False, errors='coerce')
    return pd.to_datetime(val, dayfirst=True, errors='coerce')
def format_date_val(v):
    """
    Converte qualquer objeto data/timestamp/string para 'dd/mm/yyyy'.

    FIX: usa pd.Timestamp(v) como conversor universal para cobrir pd.Timestamp,
    numpy.datetime64 e datetime.datetime — todos os tipos que df.iterrows() pode
    retornar para colunas datetime, dependendo da versão do pandas.
    """
    try:
        if pd.isna(v):
            return ''
    except (TypeError, ValueError):
        pass

    # Converte qualquer tipo date-like para Timestamp de forma segura
    try:
        ts = pd.Timestamp(v)
        if pd.isna(ts):
            return ''
        return ts.strftime('%d/%m/%Y')
    except Exception:
        pass

    # Fallback para strings
    v_str = str(v).strip()
    if v_str.lower() in ('nat', 'nan', 'none', ''):
        return ''
    try:
        if ISO_DATE.match(v_str):
            return pd.to_datetime(v_str, dayfirst=False, errors='coerce').strftime('%d/%m/%Y')
        return pd.to_datetime(v_str, dayfirst=True, errors='coerce').strftime('%d/%m/%Y')
    except Exception:
        return v_str

# ─── Regras de negócio ────────────────────────────────────────────────────────

ESTAGIO_KW   = ['estagiário', 'estagiario', 'estágio', 'estagio', 'estag']
APRENDIZ_KW  = ['aprendiz', 'menor aprendiz', 'jovem aprendiz']
TERCEIRO_KW  = ['terceiro', 'terceirizado', 'prestador']
PJ_KW        = ['pj', 'pessoa jurídica', 'pessoa juridica', 'p.j.', 'autonomo', 'autônomo']
DESCONTO_POS = {'sim', 's', 'yes', 'y', 'ativo', 'habilitado', 'enabled', 'true', '1', 'ativado', 'ok'}

def is_valid_email(v):
    return bool(EMAIL_RE.match(norm(v)))

def get_email_domain(v):
    email = norm_ascii(v)
    if '@' not in email:
        return ''
    return email.rsplit('@', 1)[-1].strip()

def is_allowed_email_domain(v):
    return get_email_domain(v) in ALLOWED_EMAIL_DOMAINS

def is_ihm_group_email_domain_allowed(domain):
    domain = norm_ascii(domain)
    return 'ihm' in domain or 'stefanini' in domain

def is_future_date(v):
    if v is None or (hasattr(v, '__class__') and v.__class__.__name__ == 'NaTType'):
        return False
    if isinstance(v, pd.Timestamp):
        return not pd.isna(v) and v.date() > date.today()
    try:
        ts = pd.Timestamp(v)
        return not pd.isna(ts) and ts.date() > date.today()
    except Exception:
        return False

def detect_vinculo_label(v):
    v = norm(v)
    if any(k in v for k in ESTAGIO_KW):  return 'Estagiário'
    if any(k in v for k in APRENDIZ_KW): return 'Aprendiz'
    if any(k in v for k in TERCEIRO_KW): return 'Terceiro/Terceirizado'
    if any(k in v for k in PJ_KW):       return 'Pessoa Jurídica (PJ)'
    return None

def auto_map_cols(headers):
    """
    Mapeia automaticamente os cabeçalhos de uma aba para os campos internos.
    Para no primeiro alias que bater — exact match tem prioridade sobre partial.
    """
    mapping = {}
    headers_lower = [(h, norm(h)) for h in headers]
    for key, aliases in COL_ALIASES.items():
        for alias in aliases:
            for h, l in headers_lower:
                if l == alias or l == alias.replace(' ', '_'):
                    mapping[key] = h
                    break
                if len(alias) >= 5 and alias in l:
                    if key not in mapping:
                        mapping[key] = h
            if key in mapping:
                break
    return mapping

# ─── Carregamento centralizado ────────────────────────────────────────────────

def _load_and_rename(xl, sheet_map, per_sheet_map):
    """
    Fonte única de verdade para leitura, limpeza e renomeação de colunas.
    Elimina a dessincronização que existia quando o mesmo bloco era duplicado
    em api_process, api_preview e api_get_columns.

    Correções aplicadas aqui:
    - Remove BOM (\\ufeff) dos cabeçalhos antes de qualquer comparação.
    - Elimina linhas completamente vazias com filtro compatível com na_filter=False.
    - Deduplica por CPF após concat para remover colaboradores que aparecem em
      mais de uma aba (609 matrículas duplicadas identificadas nos outputs).
    """
    frames = []
    for grupo, sheet_name in sheet_map.items():
        df = xl.parse(sheet_name, dtype=str, na_filter=False)

        # Remove BOM e espaços dos nomes de coluna
        df.columns = [str(c).strip().lstrip('\ufeff') for c in df.columns]

        # Remove linhas onde TODAS as células são vazias (compatível com na_filter=False)
        mask = df.apply(lambda row: row.str.strip().ne('').any(), axis=1)
        df = df[mask].reset_index(drop=True)

        sheet_cols    = list(df.columns)
        sheet_auto    = auto_map_cols(sheet_cols)
        sheet_col_map = per_sheet_map.get(grupo, {})

        rename = {}
        for key in COL_ALIASES.keys():
            user_col = sheet_col_map.get(key, '')
            if user_col and user_col in sheet_cols:
                rename[user_col] = f'_f_{key}'
            elif sheet_auto.get(key) and sheet_auto[key] in sheet_cols:
                rename[sheet_auto[key]] = f'_f_{key}'

        if norm_ascii(grupo) == 'stefanini':
            stefanini_matricula_col = next(
                (
                    col for col in sheet_cols
                    if norm_ascii(col) in {'c. custo', 'c.custo', 'c custo'}
                ),
                None
            )
            if stefanini_matricula_col:
                rename = {
                    col: target for col, target in rename.items()
                    if target != '_f_matricula'
                }
                rename[stefanini_matricula_col] = '_f_matricula'

        df = df.rename(columns=rename)
        df['_grupo'] = grupo
        frames.append(df)

    if not frames:
        return pd.DataFrame()

    consolidated = pd.concat(frames, ignore_index=True)

    # Deduplica por CPF: colaboradores que aparecem em mais de uma aba
    # (ex: Stefanini + IHM) entravam duplicados no output. CPF é usado como
    # chave porque é único por pessoa; matrícula pode repetir entre empresas.
    cpf_col = '_f_cpf'
    if cpf_col in consolidated.columns:
        consolidated['_dedupe_cpf_norm'] = consolidated[cpf_col].apply(
            lambda v: CPF_CLEAN.sub('', str(v)) if clean_val(v) else ''
        )
        consolidated['_dedupe_original_order'] = range(len(consolidated))
        consolidated = consolidated[consolidated['_dedupe_cpf_norm'].ne('')].copy()

        vinculo_col = '_f_tipo_vinculo'
        if vinculo_col in consolidated.columns:
            # Mesmo CPF em mais de uma aba: se existir uma ocorrencia nao-PJ,
            # ela tem prioridade. Se todas forem PJ, a regra remove normalmente.
            consolidated['_dedupe_pj_rank'] = consolidated[vinculo_col].apply(
                lambda v: 1 if any(k in norm_ascii(v) for k in PJ_KW) else 0
            )
            consolidated = consolidated.sort_values(
                ['_dedupe_cpf_norm', '_dedupe_pj_rank', '_dedupe_original_order']
            )
            consolidated = consolidated.drop_duplicates('_dedupe_cpf_norm', keep='first')
            consolidated = consolidated.sort_values('_dedupe_original_order')
            consolidated = consolidated.drop(columns=['_dedupe_pj_rank'])
        else:
            consolidated = consolidated.drop_duplicates('_dedupe_cpf_norm', keep='first')

        consolidated = consolidated.drop(columns=['_dedupe_cpf_norm', '_dedupe_original_order'])
        consolidated = consolidated.reset_index(drop=True)
    return consolidated


def _apply_date_conversion(df, internal_map):
    """Converte a coluna de admissão para Timestamp de forma centralizada."""
    adm_col = internal_map.get('data_admissao')
    if adm_col and adm_col in df.columns:
        df[adm_col] = df[adm_col].apply(parse_mixed_dates)
    return df, adm_col

# ─── Processamento por processo ───────────────────────────────────────────────

def apply_rules(df, col_map, process_name, check_fn):
    eligible_rows, excluded_rows = [], []
    email_col    = col_map.get('email')
    admissao_col = col_map.get('data_admissao')
    vinculo_col  = col_map.get('tipo_vinculo')

    _cols = list(df.columns)
    for _r in df.itertuples(index=False, name=None):
        row_dict = dict(zip(_cols, _r))
        email   = clean_val(row_dict.get(email_col, ''))   if email_col   else ''
        adm     = row_dict.get(admissao_col)               if admissao_col else None
        vinculo = clean_val(row_dict.get(vinculo_col, '')) if vinculo_col  else ''

        if not email or not is_valid_email(email):
            excluded_rows.append({**row_dict, '_motivo': 'E-mail ausente', '_processo': process_name})
            continue

        if adm is not None and is_future_date(adm):
            excluded_rows.append({**row_dict, '_motivo': 'Data de admissão futura', '_processo': process_name})
            continue

        reason = check_fn(row_dict, col_map, vinculo)
        if reason:
            excluded_rows.append({**row_dict, '_motivo': reason, '_processo': process_name})
            continue

        email_domain = get_email_domain(email)
        domain_allowed = email_domain in ALLOWED_EMAIL_DOMAINS

        # Regras de domínio de e-mail por grupo
        if row_dict.get('_grupo') == 'IHM' and email_domain:
            if not is_ihm_group_email_domain_allowed(email_domain):
                excluded_rows.append({**row_dict, '_motivo': 'E-mail sem domínio IHM/Stefanini', '_processo': process_name})
                continue

        TOPAZ_KW = ['topaz', 'tpz', 'grupotopaz', 'grupo-topaz', 'newm', 'top-systems', 'topsystems']
        if row_dict.get('_grupo') == 'Topaz' and email_domain and not domain_allowed:
            if not any(kw in email_domain for kw in TOPAZ_KW):
                excluded_rows.append({**row_dict, '_motivo': 'E-mail sem domínio Topaz', '_processo': process_name})
                continue

        STEFANINI_KW = ['stefanini']
        if row_dict.get('_grupo') == 'Stefanini' and email_domain and not domain_allowed:
            if not any(kw in email_domain for kw in STEFANINI_KW):
                excluded_rows.append({**row_dict, '_motivo': 'E-mail sem domínio Stefanini', '_processo': process_name})
                continue

        eligible_rows.append(row_dict)
    return eligible_rows, excluded_rows

def check_tp(row, col_map, vinculo):
    return detect_vinculo_label(vinculo) if vinculo else None

def check_wh(row, col_map, vinculo):
    """
    FIX: vinculo era comparado em case original; keywords são lowercase.
    norm() garante comparação case-insensitive independente da planilha.
    """
    vinculo_norm = norm(vinculo)
    if vinculo_norm:
        bad = ESTAGIO_KW + APRENDIZ_KW + TERCEIRO_KW + PJ_KW
        if any(k in vinculo_norm for k in bad):
            return detect_vinculo_label(vinculo) or 'Vínculo inelegível para Wellhub'
    desconto_col = col_map.get('desconto_folha')
    if desconto_col and desconto_col in row:
        val = clean_val(row.get(desconto_col))
        if val:
            if norm(val) not in DESCONTO_POS:
                return 'Sem desconto em folha habilitado'
    return None

def check_nv(row, col_map, vinculo):
    vinculo_norm = norm(vinculo)
    if vinculo_norm and any(k in vinculo_norm for k in PJ_KW):
        return 'Pessoa Jurídica (PJ)'
    return None

WELLHUB_COLS = [
    'CNPJ',
    'Nome da Empresa',
    'Nome do profissional',
    'Data de Admissão',
    'Email do profissional',
    'CPF',
    'Matrícula',
    'Department',
    'Cost Center',
    'Office Zip',
    'Code Payroll',
    'Payroll Enabled',
]

WELLHUB_EMPRESAS_KW = [
    ('gauge',), ('woopi',), ('stefanini', 'data', 'analytics'),
    ('tpz',), ('w3haus',), ('cyber', 'seguranca'),
    ('orbitall', 'atendimento'), ('orbitall', 'servicos'),
    ('orbitall', 'processamentos'), ('stefanini', 'consultoria', 'assessoria'),
    ('intelligenti',), ('n 1 software',), ('n1 software',),
]

def is_wellhub_empresa(nome):
    n = norm(nome)
    return any(all(kw in n for kw in kws) for kws in WELLHUB_EMPRESAS_KW)

# ─── Geração de Excel ─────────────────────────────────────────────────────────

FINAL_COLS = [
    'Nome da Empresa',
    'CNPJ da Empresa',
    'Matrícula',
    'Nome do profissional',
    'Função do profissional',
    'Data de Admissão',
    'Email Funcional',
    'CPF do profissional',
]
AUDIT_COLS = FINAL_COLS + ['Processo', 'Motivo da exclusão', 'Data de processamento']

COL_WIDTHS = {
    'Nome da Empresa': 35, 'CNPJ da Empresa': 22, 'Matrícula': 22,
    'Nome do profissional': 42, 'Função do profissional': 30,
    'Data de Admissão': 18, 'Email Funcional': 38, 'CPF do profissional': 22,
    'Processo': 14, 'Motivo da exclusão': 38, 'Data de processamento': 22,
}

def to_final_row(row_dict, col_map):
    get = lambda k: clean_val(row_dict.get(col_map.get(k, '___absent___'), ''))
    adm_col = col_map.get('data_admissao', '')
    adm_val = row_dict.get(adm_col, '') if adm_col else ''
    return {
        'Nome da Empresa':        get('empresa'),
        'CNPJ da Empresa':        get('cnpj'),
        'Matrícula':              get('matricula'),
        'Nome do profissional':   get('nome'),
        'Função do profissional': get('cargo'),
        'Data de Admissão':       format_date_val(adm_val),
        'Email Funcional':        get('email'),
        # FIX: format_cpf garante formatação padronizada para CPFs sem pontuação (Topaz).
        'CPF do profissional':    format_cpf(row_dict.get(col_map.get('cpf', '___absent___'), '')),
    }

def to_wellhub_row(row_dict, col_map):
    get = lambda k: clean_val(row_dict.get(col_map.get(k, '___absent___'), ''))
    adm_col = col_map.get('data_admissao', '')
    adm_val = row_dict.get(adm_col, '') if adm_col else ''
    return {
        'CNPJ':                  get('cnpj'),
        'Nome da Empresa':       get('empresa'),
        'Nome do profissional':  get('nome'),
        'Data de Admissão':      format_date_val(adm_val),
        'Email do profissional': get('email'),
        'CPF':                   format_cpf(row_dict.get(col_map.get('cpf', '___absent___'), '')),
        'Matrícula':             get('matricula'),
        'Department':            '',
        'Cost Center':           '',
        'Office Zip':            '',
        'Code Payroll':          '',
        'Payroll Enabled':       'YES',
    }


def write_fast_excel(rows, filepath, sheet_name, columns):
    """Grava Excel com pandas (bulk) + openpyxl só no cabeçalho — 10x mais rápido."""
    from openpyxl import load_workbook as _lw
    _df = pd.DataFrame(rows if rows else [dict.fromkeys(columns, '')])
    _df = _df.reindex(columns=columns, fill_value='')
    for c in _df.columns:
        _df[c] = _df[c].astype(str).replace({'nan':'','None':''})
    _df.to_excel(filepath, sheet_name=sheet_name[:31], index=False, engine='openpyxl')
    _wb = _lw(filepath)
    _ws = _wb[sheet_name[:31]]
    _hf = PatternFill('solid', fgColor='0F2A56')
    _ff = Font(name='Calibri', bold=True, color='FFFFFF', size=11)
    for ci, cn in enumerate(columns, 1):
        c = _ws.cell(row=1, column=ci)
        c.font = _ff; c.fill = _hf
        c.alignment = Alignment(horizontal='left', vertical='center')
        _ws.column_dimensions[c.column_letter].width = COL_WIDTHS.get(cn, 20)
    _ws.freeze_panes = 'A2'
    _ws.row_dimensions[1].height = 22
    _wb.save(filepath)


def write_styled_excel(rows, filepath, sheet_name, columns):
    """
    Writer rápido para bases grandes.
    Mantém cabeçalho, largura de colunas e dados.
    Remove estilização pesada célula a célula para ganhar performance.
    """
    import xlsxwriter

    wb = xlsxwriter.Workbook(filepath, {"constant_memory": True})
    ws = wb.add_worksheet(sheet_name[:31])

    header_fmt = wb.add_format({
        "bold": True,
        "font_color": "white",
        "bg_color": "#0F2A56",
        "align": "left",
        "valign": "vcenter"
    })

    for col_idx, col_name in enumerate(columns):
        ws.write(0, col_idx, col_name, header_fmt)
        ws.set_column(col_idx, col_idx, COL_WIDTHS.get(col_name, 20))

    for row_idx, row_data in enumerate(rows, start=1):
        for col_idx, col_name in enumerate(columns):
            ws.write(row_idx, col_idx, str(row_data.get(col_name, "") or ""))

    ws.freeze_panes(1, 0)
    wb.close()

# ─── Rotas Flask ──────────────────────────────────────────────────────────────


# ─── Geração de CSV por cliente ──────────────────────────────────────────────
# Colunas e mapeamentos definidos pelos layouts oficiais de cada cliente.
# Nenhuma regra de negócio é alterada aqui — só o formato de saída.

def write_totalpass_csv(tp_rows, filepath):
    """
    Layout TotalPass: CNPJ_EMPRESA | E-MAIL_COLABORADOR | MATRÍCULA
    """
    import csv
    with open(filepath, 'w', newline='', encoding='utf-8-sig') as f:
        writer = csv.DictWriter(f, fieldnames=['CNPJ_EMPRESA', 'E-MAIL_COLABORADOR', 'MATRÍCULA'], delimiter=';')
        writer.writeheader()
        for row in tp_rows:
            writer.writerow({
                'CNPJ_EMPRESA':        row.get('CNPJ da Empresa', ''),
                'E-MAIL_COLABORADOR':  row.get('Email Funcional', ''),
                'MATRÍCULA':           row.get('Matrícula', ''),
            })

def write_newvalue_csv(nv_rows, filepath):
    """
    Layout New Value: Nome | Nome Empresa | CNPJ Empresa | Número de CPF
    """
    import csv
    with open(filepath, 'w', newline='', encoding='utf-8-sig') as f:
        writer = csv.DictWriter(f, fieldnames=['Nome', 'Nome Empresa', 'CNPJ Empresa', 'Número de CPF'], delimiter=';')
        writer.writeheader()
        for row in nv_rows:
            writer.writerow({
                'Nome':           row.get('Nome do profissional', ''),
                'Nome Empresa':   row.get('Nome da Empresa', ''),
                'CNPJ Empresa':   row.get('CNPJ da Empresa', ''),
                'Número de CPF':  row.get('CPF do profissional', ''),
            })

def _write_wellhub_csv_legacy(wh_rows, filepath):
    """
    Layout Wellhub — 10 colunas com nomes exatos do sistema Wellhub.
    """
    import csv
    COLS = [
        'Name (nome) obrigatório',
        'Email obrigatório',
        'National ID (cpf) obrigatório',
        'Employee ID (matrícula) obrigatório',
        'Department',
        'Cost Center',
        'Office Zip Code',
        'Payroll ID',
        'Payroll Enabled (folha de pagamento habilitada?) obrigatório',
        'Employee Segment',
    ]
    with open(filepath, 'w', newline='', encoding='utf-8-sig') as f:
        writer = csv.DictWriter(f, fieldnames=COLS, delimiter=';')
        writer.writeheader()
        for row in wh_rows:
            writer.writerow({
                'Name (nome) obrigatório':                                   row.get('Name', ''),
                'Email obrigatório':                                         row.get('Email', ''),
                'National ID (cpf) obrigatório':                             row.get('National ID', ''),
                'Employee ID (matrícula) obrigatório':                       row.get('Employee ID', ''),
                'Department':                                                row.get('Department', ''),
                'Cost Center':                                               row.get('Cost Center', ''),
                'Office Zip Code':                                           row.get('Office Zip Code', ''),
                'Payroll ID':                                                row.get('Payroll ID', ''),
                'Payroll Enabled (folha de pagamento habilitada?) obrigatório': row.get('Payroll Enabled', 'YES'),
                'Employee Segment':                                          '',
            })

def write_wellhub_csv(wh_rows, filepath):
    """Layout WellHub consolidado para envio ao cliente."""
    import csv
    cols = [
        'CNPJ',
        'Nome da Empresa',
        'Nome do profissional',
        'Data de Admissão',
        'Email do profissional',
        'CPF',
        'Matrícula',
        'Department',
        'Cost Center',
        'Office Zip',
        'Code Payroll',
        'Payroll Enabled',
    ]
    with open(filepath, 'w', newline='', encoding='utf-8-sig') as f:
        writer = csv.DictWriter(f, fieldnames=cols, delimiter=';')
        writer.writeheader()
        for row in wh_rows:
            writer.writerow({
                'CNPJ': row.get('CNPJ', ''),
                'Nome da Empresa': row.get('Nome da Empresa', ''),
                'Nome do profissional': row.get('Nome do profissional', ''),
                'Data de Admissão': row.get('Data de Admissão', ''),
                'Email do profissional': row.get('Email do profissional', ''),
                'CPF': row.get('CPF', ''),
                'Matrícula': row.get('Matrícula', ''),
                'Department': row.get('Department', ''),
                'Cost Center': row.get('Cost Center', ''),
                'Office Zip': row.get('Office Zip', ''),
                'Code Payroll': row.get('Code Payroll', ''),
                'Payroll Enabled': row.get('Payroll Enabled', 'YES'),
            })

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/api/upload', methods=['POST'])
def api_upload():
    if 'file' not in request.files:
        return jsonify({'error': 'Nenhum arquivo enviado.'}), 400
    f = request.files['file']
    if not f.filename.lower().endswith(('.xlsx', '.xls')):
        return jsonify({'error': 'Formato inválido. Envie um arquivo .xlsx ou .xls.'}), 400

    session_id  = str(uuid.uuid4())
    session_dir = os.path.join(TEMP_DIR, session_id)
    os.makedirs(session_dir, exist_ok=True)
    upload_path = os.path.join(session_dir, 'source.xlsx')
    f.save(upload_path)

    try:
        xl       = pd.ExcelFile(upload_path)
        auto_map = {}
        for grupo in ['Stefanini', 'Topaz', 'IHM']:
            for s in xl.sheet_names:
                if grupo.lower() in s.lower() or s.lower() in grupo.lower():
                    auto_map[grupo] = s
                    break
        return jsonify({
            'session_id':     session_id,
            'sheet_names':    xl.sheet_names,
            'auto_sheet_map': auto_map,
            'filename':       f.filename,
        })
    except Exception as e:
        return jsonify({'error': f'Erro ao ler o arquivo: {str(e)}'}), 500

@app.route('/api/get-columns', methods=['POST'])
def api_get_columns():
    data        = request.json
    sid         = data.get('session_id')
    sheet_map   = data.get('sheet_map', {})
    upload_path = os.path.join(TEMP_DIR, sid, 'source.xlsx')
    if not os.path.exists(upload_path):
        return jsonify({'error': 'Sessão expirada.'}), 400

    try:
        xl         = pd.ExcelFile(upload_path)
        sheet_data = {}
        for grupo, sheet_name in sheet_map.items():
            df = xl.parse(sheet_name, dtype=str, na_filter=False)
            df.columns = [str(c).strip().lstrip('\ufeff') for c in df.columns]
            mask = df.apply(lambda row: row.str.strip().ne('').any(), axis=1)
            df   = df[mask].reset_index(drop=True)
            cols = list(df.columns)
            sheet_data[grupo] = {
                'headers':   cols,
                'auto_map':  auto_map_cols(cols),
                'row_count': len(df),
            }
        return jsonify({'sheet_data': sheet_data})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/process', methods=['POST'])
def api_process():
    data          = request.json
    sid           = data.get('session_id')
    sheet_map     = data.get('sheet_map', {})
    per_sheet_map = data.get('per_sheet_col_map', {})
    upload_path   = os.path.join(TEMP_DIR, sid, 'source.xlsx')
    if not os.path.exists(upload_path):
        return jsonify({'error': 'Sessão expirada.'}), 400

    try:
        xl           = pd.ExcelFile(upload_path)
        consolidated = _load_and_rename(xl, sheet_map, per_sheet_map)
        internal_map = {key: f'_f_{key}' for key in COL_ALIASES.keys()}

        consolidated, _ = _apply_date_conversion(consolidated, internal_map)

        today_str   = date.today().strftime('%d/%m/%Y')
        session_dir = os.path.join(TEMP_DIR, sid)

        tp_elig, tp_excl = apply_rules(consolidated, internal_map, 'TotalPass', check_tp)
        wh_elig, wh_excl = apply_rules(consolidated, internal_map, 'Wellhub',   check_wh)
        nv_elig, nv_excl = apply_rules(consolidated, internal_map, 'New Value', check_nv)

        tp_rows = [to_final_row(r, internal_map) for r in tp_elig]
        wh_rows = [to_final_row(r, internal_map) for r in wh_elig]
        nv_rows = [to_final_row(r, internal_map) for r in nv_elig]

        audit_rows = []
        for excl_list in (tp_excl, wh_excl, nv_excl):
            for r in excl_list:
                row = to_final_row(r, internal_map)
                row['Processo']              = r.get('_processo', '')
                row['Motivo da exclusão']    = r.get('_motivo', '')
                row['Data de processamento'] = today_str
                audit_rows.append(row)

        write_styled_excel(tp_rows,    os.path.join(session_dir, 'totalpass.xlsx'), 'TotalPass', FINAL_COLS)
        write_styled_excel(nv_rows,    os.path.join(session_dir, 'newvalue.xlsx'),  'New Value',  FINAL_COLS)
        write_styled_excel(audit_rows, os.path.join(session_dir, 'auditoria.xlsx'), 'Auditoria',  AUDIT_COLS)

        # ── Gera CSVs no formato de cada cliente ──────────────────────────
        write_totalpass_csv(tp_rows, os.path.join(session_dir, 'totalpass.csv'))
        write_newvalue_csv(nv_rows,  os.path.join(session_dir, 'newvalue.csv'))

        # ──────────────────────────────────────────────────────────────────

        # Salva total da base para o PPT de auditoria
        with open(os.path.join(session_dir, 'total_base.txt'), 'w') as _f:
            _f.write(str(len(consolidated)))

        wh_by_cnpj = {}
        cnpj_col   = internal_map.get('cnpj', '')
        emp_col    = internal_map.get('empresa', '')
        for r in wh_elig:
            cnpj_raw = clean_val(r.get(cnpj_col, ''))
            emp_raw  = clean_val(r.get(emp_col, ''))
            cnpj     = cnpj_raw if cnpj_raw else f"GRUPO_{r.get('_grupo', '')}"
            empresa  = emp_raw  if emp_raw  else r.get('_grupo', '')
            if cnpj not in wh_by_cnpj:
                wh_by_cnpj[cnpj] = {'empresa': empresa, 'rows': []}
            wh_row = to_wellhub_row(r, internal_map)
            wh_by_cnpj[cnpj]['rows'].append(wh_row)

        wh_by_cnpj   = {c: g for c, g in wh_by_cnpj.items() if is_wellhub_empresa(g['empresa'])}
        wh_file_meta = []
        for cnpj, grp in wh_by_cnpj.items():
            safe_name = re.sub(r'[^\w]', '_', grp['empresa']).upper().strip('_')
            fname     = f"{safe_name}_Wellhub.xlsx"
            fpath     = os.path.join(session_dir, f"wh_{re.sub(r'[^\w]', '', cnpj)}.xlsx")
            write_styled_excel(grp['rows'], fpath, 'Wellhub', WELLHUB_COLS)
            # guarda rows para uso no CSV
            _rows_raw = grp['rows'][:]
            wh_file_meta.append({
                'cnpj':     cnpj,
                'empresa':  grp['empresa'],
                'count':    len(grp['rows']),
                'filepath': fpath,
                'filename': fname,
                '_rows_raw': _rows_raw,
            })

        zip_path = os.path.join(session_dir, 'wellhub_todos.zip')
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for meta in wh_file_meta:
                zf.write(meta['filepath'], meta['filename'])

        with open(os.path.join(session_dir, 'wh_meta.json'), 'w', encoding='utf-8') as jf:
            json.dump(wh_file_meta, jf, ensure_ascii=False)

        wh_csv_rows = []
        for meta in wh_file_meta:
            wh_csv_rows.extend(meta.get('_rows_raw', []))
        write_wellhub_csv(wh_csv_rows, os.path.join(session_dir, 'wellhub_geral.csv'))

        # Mantem um ZIP legado com o CSV geral para compatibilidade interna.
        import csv as _csv, zipfile as _zf
        _WCOLS = ['Name (nome) obrigatório','Email obrigatório',
                  'National ID (cpf) obrigatório','Employee ID (matrícula) obrigatório',
                  'Department','Cost Center','Office Zip Code','Payroll ID',
                  'Payroll Enabled (folha de pagamento habilitada?) obrigatório','Employee Segment']
        _wh_csvs = [(os.path.join(session_dir, 'wellhub_geral.csv'), 'Wellhub_Geral.csv')]
        for meta in []:
            _safe     = re.sub(r'[^\w]', '_', meta['empresa']).upper().strip('_')
            _csv_path = os.path.join(session_dir, f"wh_{re.sub(r'[^\w]','',meta['cnpj'])}.csv")
            _csv_name = f"{_safe}_Wellhub.csv"
            meta['csv_path'] = _csv_path
            meta['csv_name'] = _csv_name
            import pandas as _pd
            _df = _pd.read_excel(meta['filepath'], dtype=str)
            with open(_csv_path, 'w', newline='', encoding='utf-8-sig') as _f:
                _w = _csv.writer(_f, delimiter=';')
                _w.writerow(_WCOLS)
                for _, _row in _df.iterrows():
                    _w.writerow([
                        _row.get('Name',''), _row.get('Email',''),
                        _row.get('National ID',''), _row.get('Employee ID',''),
                        '','','','', _row.get('Payroll Enabled','YES'), ''
                    ])
            _wh_csvs.append((_csv_path, _csv_name))
        _wh_zip = os.path.join(session_dir, 'wellhub_csv_todos.zip')
        with _zf.ZipFile(_wh_zip, 'w', _zf.ZIP_DEFLATED) as _z:
            for _p, _n in _wh_csvs:
                _z.write(_p, _n)

        with open(os.path.join(session_dir, 'wh_meta.json'), 'w', encoding='utf-8') as jf:
            json.dump(wh_file_meta, jf, ensure_ascii=False)


        def count_reasons(rows):
            counts = {}
            for r in rows:
                key = r.get('_motivo', 'Desconhecido')
                counts[key] = counts.get(key, 0) + 1
            return counts

        return jsonify({
            'total': len(consolidated),
            'tp': {'eligible': len(tp_rows),  'excluded': len(tp_excl),  'reasons': count_reasons(tp_excl)},
            'wh': {
                'eligible':  sum(m['count'] for m in wh_file_meta),
                'excluded':  len(wh_excl),
                'reasons':   count_reasons(wh_excl),
                'companies': wh_file_meta,
            },
            'nv': {'eligible': len(nv_rows),  'excluded': len(nv_excl),  'reasons': count_reasons(nv_excl)},
            'audit_count': len(audit_rows),
        })
    except Exception as e:
        import traceback
        return jsonify({'error': str(e), 'detail': traceback.format_exc()}), 500

@app.route('/api/preview', methods=['POST'])
def api_preview():
    data          = request.json
    sid           = data.get('session_id')
    sheet_map     = data.get('sheet_map', {})
    per_sheet_map = data.get('per_sheet_col_map', {})
    upload_path   = os.path.join(TEMP_DIR, sid, 'source.xlsx')
    if not os.path.exists(upload_path):
        return jsonify({'error': 'Sessão expirada.'}), 400

    try:
        xl           = pd.ExcelFile(upload_path)
        consolidated = _load_and_rename(xl, sheet_map, per_sheet_map)
        internal_map = {key: f'_f_{key}' for key in COL_ALIASES.keys()}
        consolidated, _ = _apply_date_conversion(consolidated, internal_map)

        def cell(row, key):
            return clean_val(row.get(internal_map.get(key, '___absent___'), ''))

        rows = []
        for _, row in consolidated.head(15).iterrows():
            rows.append({
                'empresa':      cell(row, 'empresa'),
                'nome':         cell(row, 'nome'),
                'matricula':    cell(row, 'matricula'),
                'email':        cell(row, 'email'),
                'tipo_vinculo': cell(row, 'tipo_vinculo'),
                'cargo':        cell(row, 'cargo'),
                '_grupo':       row.get('_grupo', ''),
            })

        return jsonify({'rows': rows, 'total': len(consolidated)})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/wellhub-zip/<session_id>', methods=['POST'])
def wellhub_zip_filtered(session_id):
    session_dir = os.path.join(TEMP_DIR, session_id)
    meta_path   = os.path.join(session_dir, 'wh_meta.json')
    if not os.path.exists(meta_path):
        return jsonify({'error': 'Sessão expirada.'}), 400

    selected_cnpjs = set(request.json.get('cnpjs', []))
    with open(meta_path, encoding='utf-8') as f:
        wh_file_meta = json.load(f)

    selected = [m for m in wh_file_meta if m['cnpj'] in selected_cnpjs]
    if not selected:
        return jsonify({'error': 'Nenhuma empresa selecionada.'}), 400

    zip_path = os.path.join(session_dir, 'wellhub_filtrado.zip')
    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
        for meta in selected:
            if os.path.exists(meta['filepath']):
                zf.write(meta['filepath'], meta['filename'])

    return send_file(zip_path, as_attachment=True, download_name=f'Wellhub_{len(selected)}_Empresas.zip')

@app.route('/api/download/<session_id>/<file_type>')
def api_download(session_id, file_type):
    session_dir = os.path.join(TEMP_DIR, session_id)
    if file_type == 'wellhub':
        path = os.path.join(session_dir, 'wellhub_todos.zip')
        if os.path.exists(path):
            return send_file(path, as_attachment=True, download_name='Wellhub_Todos_CNPJs.zip')
        return 'Arquivo não encontrado.', 404

    file_map = {
        'totalpass': ('totalpass.xlsx', 'Base_TotalPass_Tratada.xlsx'),
        'newvalue':  ('newvalue.xlsx',  'Base_NewValue_Tratada.xlsx'),
        'auditoria': ('auditoria.xlsx', 'Base_Auditoria_Exclusoes.xlsx'),
    }
    if file_type not in file_map:
        return 'Tipo inválido.', 404
    internal, display = file_map[file_type]
    path = os.path.join(session_dir, internal)
    if os.path.exists(path):
        return send_file(path, as_attachment=True, download_name=display)
    return 'Arquivo não encontrado.', 404



@app.route('/api/download/<session_id>/auditppt')
def download_audit_ppt(session_id):
    session_dir = os.path.join(TEMP_DIR, session_id)
    audit_path  = os.path.join(session_dir, 'auditoria.xlsx')
    ppt_path    = os.path.join(session_dir, 'auditoria_ppt.pptx')
    if not os.path.exists(audit_path):
        return 'Sessão expirada ou processamento não concluído.', 404
    try:
        total_base_path = os.path.join(session_dir, 'total_base.txt')
        total_base = int(open(total_base_path).read().strip()) if os.path.exists(total_base_path) else 0
        generate_ppt(audit_path, ppt_path, total_base)
        return send_file(ppt_path, as_attachment=True,
                         download_name='Auditoria_Exclusoes_BenefProcess.pptx')
    except Exception as e:
        import traceback
        return f'Erro ao gerar PPT: {str(e)}\n{traceback.format_exc()}', 500


@app.route('/api/download/<session_id>/totalpass_csv')
def download_totalpass_csv(session_id):
    path = os.path.join(TEMP_DIR, session_id, 'totalpass.csv')
    if not os.path.exists(path): return 'Arquivo não encontrado.', 404
    return send_file(path, as_attachment=True, download_name='TotalPass_Colaboradores.csv')

@app.route('/api/download/<session_id>/newvalue_csv')
def download_newvalue_csv(session_id):
    path = os.path.join(TEMP_DIR, session_id, 'newvalue.csv')
    if not os.path.exists(path): return 'Arquivo não encontrado.', 404
    return send_file(path, as_attachment=True, download_name='NewValue_Colaboradores.csv')

@app.route('/api/download/<session_id>/wellhub_csv_zip')
def download_wellhub_csv_zip(session_id):
    path = os.path.join(TEMP_DIR, session_id, 'wellhub_geral.csv')
    if not os.path.exists(path): return 'Arquivo não encontrado.', 404
    return send_file(path, as_attachment=True, download_name='Wellhub_Geral.csv')

@app.route('/api/wellhub-csv-zip-filtered/<session_id>', methods=['POST'])
def wellhub_csv_zip_filtered(session_id):
    """Download ZIP com CSVs Wellhub apenas das empresas selecionadas."""
    session_dir    = os.path.join(TEMP_DIR, session_id)
    meta_path      = os.path.join(session_dir, 'wh_meta.json')
    if not os.path.exists(meta_path): return jsonify({'error': 'Sessão expirada.'}), 400
    selected_cnpjs = set(request.json.get('cnpjs', []))
    with open(meta_path, encoding='utf-8') as f: wh_file_meta = json.load(f)
    selected = [m for m in wh_file_meta if m['cnpj'] in selected_cnpjs]
    if not selected: return jsonify({'error': 'Nenhuma empresa selecionada.'}), 400
    zip_path = os.path.join(session_dir, 'wellhub_csv_filtrado.zip')
    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
        for meta in selected:
            csv_p = meta.get('csv_path', '')
            if os.path.exists(csv_p):
                zf.write(csv_p, meta.get('csv_name', os.path.basename(csv_p)))
    return send_file(zip_path, as_attachment=True,
                     download_name=f'Wellhub_CSV_{len(selected)}_Empresas.zip')


# ══════════════════════════════════════════════════════════════════════════════
# MAR SAÚDE — fluxo independente, sem alterar nenhuma regra existente
# ══════════════════════════════════════════════════════════════════════════════

MAR_SAUDE_COLS = [
    'Empresa', 'Nome Empresa', 'CNPJ Empresa', 'Célula',
    'Desc. Atividade (Serviço)', 'Nome', 'Nome Social', 'Nome Cargo',
    'Data de Admissão', 'Sexo', 'Data de Nascimento', 'Logradouro',
    'Endereço', 'Complemento', 'Bairro', 'Cidade', 'Complemento CEP',
    'E-mail Funcional', 'Número de CPF', 'Desc. Unidade Adm. (Cliente)',
    'Desc. Tipo de Vínculo',
]

MAR_SAUDE_COL_ALIASES = {
    'empresa':       ['nome empresa', 'empresa'],
    'nome_empresa':  ['nome empresa', 'empresa'],
    'cnpj':          ['cnpj empresa', 'cnpj da empresa', 'cnpj'],
    'celula':        ['c. custo', 'célula', 'celula', 'centro de custo'],
    'atividade':     ['desc. atividade (serviço)', 'desc. atividade', 'atividade (serviço)', 'atividade'],
    'nome':          ['nome', 'nome profissional', 'nome do profissional'],
    'nome_social':   ['nome social'],
    'cargo':         ['nome cargo', 'cargo', 'nome função', 'nome funcao', 'função'],
    'data_admissao': ['data de admissão', 'data de admissao', 'admissão', 'admissao', 'dt admissão'],
    'sexo':          ['sexo', 'gênero', 'genero'],
    'dt_nascimento': ['data de nascimento', 'nascimento', 'dt nascimento', 'data nascimento'],
    'logradouro':    ['logradouro', 'tipo logradouro'],
    'endereco':      ['endereço', 'endereco', 'rua'],
    'complemento':   ['complemento'],
    'bairro':        ['bairro'],
    'cidade':        ['cidade', 'município', 'municipio'],
    'cep':           ['complemento cep', 'cep', 'cod postal'],
    'email':         ['e-mail funcional', 'email funcional', 'e-mail', 'email'],
    'cpf':           ['número de cpf', 'cpf', 'cpf do profissional', 'nr cpf'],
    'unidade_adm':   ['desc. unidade adm. (cliente)', 'desc. unidade adm', 'unidade adm'],
    'tipo_vinculo':  ['desc. tipo de vínculo', 'desc. tipo de vinculo',
                      'desc. vínculo', 'desc. vinculo', 'vínculo', 'vinculo', 'regime'],
    'nome_sindicato':['nome sindicato', 'sindicato', 'nome_sindicato',
                      'desc. sindicato', 'desc sindicato'],
}

def _auto_map_ms(headers):
    """Mapeamento automático usando aliases do Mar Saúde."""
    mapping = {}
    headers_lower = [(h, norm(h)) for h in headers]
    for key, aliases in MAR_SAUDE_COL_ALIASES.items():
        for alias in aliases:
            for h, l in headers_lower:
                if l == alias or l == alias.replace(' ', '_'):
                    mapping[key] = h; break
                if len(alias) >= 5 and alias in l:
                    if key not in mapping: mapping[key] = h
            if key in mapping: break
    return mapping

def _to_ms_row(row_dict, col_map):
    """Converte linha para as 20 colunas do Mar Saúde."""
    get = lambda k: clean_val(row_dict.get(col_map.get(k, '___ms___'), ''))
    def _date(k):
        col = col_map.get(k, '')
        val = row_dict.get(col, '') if col else ''
        try:
            return format_date_val(parse_mixed_dates(val)) if val else ''
        except Exception:
            return clean_val(val)
    return {
        'Empresa':                      get('empresa'),
        'Nome Empresa':                 get('nome_empresa'),
        'CNPJ Empresa':                 get('cnpj'),
        'Célula':                       get('celula'),
        'Desc. Atividade (Serviço)':    get('atividade'),
        'Nome':                         get('nome'),
        'Nome Social':                  get('nome_social'),
        'Nome Cargo':                   get('cargo'),
        'Data de Admissão':             _date('data_admissao'),
        'Sexo':                         get('sexo'),
        'Data de Nascimento':           _date('dt_nascimento'),
        'Logradouro':                   get('logradouro'),
        'Endereço':                     get('endereco'),
        'Complemento':                  get('complemento'),
        'Bairro':                       get('bairro'),
        'Cidade':                       get('cidade'),
        'Complemento CEP':              get('cep'),
        'E-mail Funcional':             get('email'),
        'Número de CPF':                get('cpf'),
        'Desc. Unidade Adm. (Cliente)': get('unidade_adm'),
        'Desc. Tipo de Vínculo':         get('tipo_vinculo'),
    }

def _write_ms_csv(rows, filepath):
    import csv
    with open(filepath, 'w', newline='', encoding='utf-8-sig') as f:
        w = csv.DictWriter(f, fieldnames=MAR_SAUDE_COLS, delimiter=';')
        w.writeheader()
        w.writerows(rows)

@app.route('/api/upload-mar-saude', methods=['POST'])
def api_upload_mar_saude():
    if 'file' not in request.files:
        return jsonify({'error': 'Nenhum arquivo enviado.'}), 400
    f = request.files['file']
    if not f.filename.lower().endswith(('.xlsx', '.xls')):
        return jsonify({'error': 'Envie um arquivo .xlsx ou .xls.'}), 400
    sid         = str(uuid.uuid4())
    session_dir = os.path.join(TEMP_DIR, sid)
    os.makedirs(session_dir, exist_ok=True)
    upload_path = os.path.join(session_dir, 'ms_source.xlsx')
    f.save(upload_path)
    try:
        xl = pd.ExcelFile(upload_path)
        return jsonify({'session_id': sid, 'sheet_names': xl.sheet_names, 'filename': f.filename})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/get-columns-mar-saude', methods=['POST'])
def api_get_columns_mar_saude():
    """Retorna somente cabeçalhos e auto-mapeamento. Não processa a planilha."""
    data        = request.json or {}
    sid         = data.get('session_id')
    sheet_name  = data.get('sheet_name')
    upload_path = os.path.join(TEMP_DIR, sid, 'ms_source.xlsx')
    if not os.path.exists(upload_path):
        return jsonify({'error': 'Sessão expirada.'}), 400
    try:
        df = pd.read_excel(upload_path, sheet_name=sheet_name, dtype=str,
                           keep_default_na=False, nrows=0)
        df.columns = [str(c).strip().lstrip('\ufeff') for c in df.columns]
        cols = list(df.columns)
        return jsonify({'headers': cols, 'auto_map': _auto_map_ms(cols), 'row_count': None})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/process-mar-saude', methods=['POST'])
def api_process_mar_saude():
    """Processa Mar Saúde. Regra: remove PJ por Nome Sindicato ou Desc. Tipo de Vínculo."""
    data        = request.json or {}
    sid         = data.get('session_id')
    sheet_name  = data.get('sheet_name')
    col_map_usr = data.get('col_map', {}) or {}
    upload_path = os.path.join(TEMP_DIR, sid, 'ms_source.xlsx')
    if not os.path.exists(upload_path):
        return jsonify({'error': 'Sessão expirada.'}), 400
    try:
        # Lê headers primeiro → seleciona só colunas necessárias
        _hdr_ms = pd.read_excel(upload_path, sheet_name=sheet_name, dtype=str,
                                keep_default_na=False, nrows=0)
        _hdr_ms.columns = [str(c).strip().lstrip('\ufeff') for c in _hdr_ms.columns]
        _auto_ms   = _auto_map_ms(list(_hdr_ms.columns))
        _fmap_peek = {**_auto_ms, **{k: v for k, v in col_map_usr.items() if v}}
        _needed_ms = list({v for v in _fmap_peek.values() if v and v in _hdr_ms.columns})
        df = pd.read_excel(upload_path, sheet_name=sheet_name, dtype=str,
                           keep_default_na=False,
                           usecols=_needed_ms if _needed_ms else None)
        df.columns = [str(c).strip().lstrip('\ufeff') for c in df.columns]
        mask = df.apply(lambda row: row.astype(str).str.strip().ne('').any(), axis=1)
        df = df[mask].reset_index(drop=True)

        auto = _auto_map_ms(list(df.columns))
        final_map = {**auto, **{k: v for k, v in col_map_usr.items() if v}}

        for dk in ('data_admissao', 'dt_nascimento'):
            col = final_map.get(dk)
            if col and col in df.columns:
                df[col] = df[col].apply(parse_mixed_dates)

        sindicato_col = final_map.get('nome_sindicato')
        vinculo_col   = final_map.get('tipo_vinculo')
        # ── Vetorizado ──
        _sind = (df[sindicato_col].str.lower().str.strip().fillna('')
                 if sindicato_col and sindicato_col in df.columns
                 else pd.Series(['']*len(df)))
        _vinc = (df[vinculo_col].str.lower().str.strip().fillna('')
                 if vinculo_col and vinculo_col in df.columns
                 else pd.Series(['']*len(df)))
        pj_ms = (_sind.str.contains('somente pj', na=False) |
                 (_vinc == 'outros') |
                 _vinc.str.contains('|'.join(PJ_KW), na=False))
        _excl = df[pj_ms].to_dict('records')
        for r in _excl: r['_motivo'] = 'Pessoa Jurídica (PJ)'
        excluded = _excl
        eligible = df[~pj_ms].to_dict('records')

        ms_rows = [_to_ms_row(r, final_map) for r in eligible]
        ex_rows = [_to_ms_row(r, final_map) for r in excluded]

        session_dir = os.path.join(TEMP_DIR, sid)
        write_styled_excel(ms_rows, os.path.join(session_dir, 'ms_output.xlsx'),
                           'Mar Saúde', MAR_SAUDE_COLS)
        _write_ms_csv(ms_rows, os.path.join(session_dir, 'ms_output.csv'))
        write_styled_excel(ex_rows, os.path.join(session_dir, 'ms_audit.xlsx'),
                           'Excluídos', MAR_SAUDE_COLS)

        return jsonify({
            'total': len(df),
            'eligible': len(ms_rows),
            'excluded': len(ex_rows),
        })
    except Exception as e:
        import traceback
        return jsonify({'error': str(e), 'detail': traceback.format_exc()}), 500

@app.route('/api/download/<session_id>/ms_excel')
def download_ms_excel(session_id):
    path = os.path.join(TEMP_DIR, session_id, 'ms_output.xlsx')
    if os.path.exists(path):
        return send_file(path, as_attachment=True, download_name='MarSaude_Elegiveis.xlsx')
    return 'Arquivo não encontrado.', 404

@app.route('/api/download/<session_id>/ms_csv')
def download_ms_csv(session_id):
    path = os.path.join(TEMP_DIR, session_id, 'ms_output.csv')
    if os.path.exists(path):
        return send_file(path, as_attachment=True, download_name='MarSaude_Elegiveis.csv')
    return 'Arquivo não encontrado.', 404

@app.route('/api/download/<session_id>/ms_audit')
def download_ms_audit(session_id):
    path = os.path.join(TEMP_DIR, session_id, 'ms_audit.xlsx')
    if os.path.exists(path):
        return send_file(path, as_attachment=True, download_name='MarSaude_Excluidos.xlsx')
    return 'Arquivo não encontrado.', 404


# ══════════════════════════════════════════════════════════════════════════════
# EXPERT
# ══════════════════════════════════════════════════════════════════════════════

EXPERT_COLS = [
    'Nome Unidade', 'Nome Setor', 'Nome Cargo', 'Matrícula',
    'Nome Funcionário', 'Dt.Nascimento', 'Sexo', 'Situação',
    'Dt.Admissão', 'CPF', 'CBO', 'Nome da Mae do Funcionário',
    'Razão Social', 'CNPJ', 'Código Categoria (eSocial)',
]

EXPERT_INATIVOS_COLS = [
    'Nome Funcionário', 'Situação Original', 'Motivo Inativação', 'Dt.Desligamento',
    'Matrícula', 'CPF', 'Dt.Nascimento', 'Sexo', 'Dt.Admissão',
    'Nome Unidade', 'Nome Setor', 'Nome Cargo', 'CBO',
    'Nome da Mae do Funcionário', 'Razão Social', 'CNPJ',
    'Código Categoria (eSocial)',
]

# Expert: a base de ativos agora e enviada separada da base de desligados.
# As situacoes abaixo, quando vierem na base de ativos, permanecem como ativas.
EXPERT_SITUACAO_ATIVA    = 'ativo'
EXPERT_SITUACOES_ATIVAS  = {
    'ativo',
    'licenca gestacao',
    'afastado inss ate 15',
    'licen.nao remunerada',
    'outros-proc judicial',
    'processo de abandono',
    'afastamento inss >15',
    'acid. trabalho > 15',
    'aposent. invalidez',
}
SITUACAO_INATIVO_CODES   = {'90','91','92','93','95','96','99','9a','9A'}

EXPERT_COL_ALIASES = {
    'nome_unidade':  ['nome filial', 'nome unidade', 'unidade', 'nome da unidade', 'filial'],
    'nome_setor':    ['desc. atividade (serviço)', 'desc. atividade', 'atividade (serviço)',
                      'nome setor', 'setor', 'departamento'],
    'nome_cargo':    ['nome cargo', 'cargo', 'função', 'nome função', 'funcao'],
    'matricula':     ['matrícula', 'matricula', 'chapa', 'matrícula rh', 'matricula rh',
                      'matricula esocial'],
    'nome':          ['nome', 'nome funcionário', 'nome funcionario', 'nome profissional'],
    'dt_nascimento': ['data de nascimento', 'dt.nascimento', 'dt nascimento', 'data nascimento'],
    'sexo':          ['sexo', 'gênero', 'genero'],
    'situacao':      ['desc. situação', 'desc. situacao', 'desc situação', 'desc situacao',
                      'descrição situação', 'descricao situacao', 'situação', 'situacao',
                      'desc. situacao funcionário', 'desc. situacao funcionario',
                      'status', 'sit.'],
    'dt_admissao':   ['data de admissão', 'data de admissao', 'dt.admissão', 'dt.admissao',
                      'dt admissão', 'dt admissao', 'admissão'],
    'dt_desligamento': ['data de desligamento', 'dt.desligamento', 'dt desligamento',
                        'data desligamento', 'desligamento', 'data demissão',
                        'data demissao', 'dt.demissão', 'dt.demissao', 'dt demissão',
                        'dt demissao', 'demissão', 'demissao', 'data de rescisão',
                        'data de rescisao', 'data rescisão', 'data rescisao',
                        'dt.rescisão', 'dt.rescisao', 'dt rescisão', 'dt rescisao'],
    'cpf':           ['número de cpf', 'numero de cpf', 'cpf', 'nr cpf', 'pis/pasep'],
    'cbo':           ['id. cbo', 'cbo', 'cbo 2002', 'cod cbo', 'código cbo', 'cód. cbo'],
    'nome_mae':      ['nome mãe', 'nome mae', 'nome da mae do funcionário', 'nome da mae'],
    'razao_social':  ['nome empresa', 'razão social', 'razao social', 'razão social unid.'],
    'cnpj':          ['cnpj empresa', 'cnpj', 'cnpj unidade', 'cnpj da empresa'],
    'cod_categoria': ['desc. tipo de vínculo', 'desc. tipo de vinculo', 'cód. categoria', 'cod. categoria', 'código categoria (esocial)',
                      'codigo categoria esocial', 'categoria esocial'],
    'tipo_vinculo':   ['desc. tipo de vínculo', 'desc. tipo de vinculo',
                       'desc. vínculo', 'desc. vinculo', 'vínculo', 'vinculo',
                       'tipo de vínculo', 'tipo de vinculo', 'regime', 'modalidade'],
    'nome_sindicato': ['nome sindicato', 'sindicato', 'nome_sindicato',
                       'desc. sindicato', 'desc sindicato'],
}

INATIVO_KW_EXPERT = ['inativo', 'desligado', 'demitido', 'rescindido', 'inactive']


def _is_pj_expert(row_d, col_map):
    """
    Regra PJ do Expert.

    Mantém o critério por Código Categoria (eSocial) e acrescenta a regra do
    Mar Saúde:
    - Nome Sindicato contendo "Somente PJ";
    - Desc. Tipo de Vínculo igual a "OUTROS";
    - Desc. Tipo de Vínculo contendo termos de PJ.
    """
    cod_cat_col = col_map.get('cod_categoria', '')
    if cod_cat_col:
        cod_cat = norm(clean_val(row_d.get(cod_cat_col, '')))
        if any(k in cod_cat for k in PJ_KW):
            return True

    sindicato_col = col_map.get('nome_sindicato', '')
    if sindicato_col:
        sindicato = norm(clean_val(row_d.get(sindicato_col, '')))
        if 'somente pj' in sindicato:
            return True

    vinculo_col = col_map.get('tipo_vinculo', '')
    if vinculo_col:
        vinculo = norm(clean_val(row_d.get(vinculo_col, '')))
        if vinculo == 'outros' or any(k in vinculo for k in PJ_KW):
            return True

    return False

COL_WIDTHS.update({
    'Nome Unidade': 35, 'Nome Setor': 30, 'Nome Cargo': 32,
    'Nome Funcionário': 42, 'Dt.Nascimento': 16, 'Sexo': 10,
    'Situação': 14, 'Dt.Admissão': 16, 'CBO': 12,
    'Nome da Mae do Funcionário': 42, 'Razão Social': 40,
    'Código Categoria (eSocial)': 30,
    'Situação Original': 26, 'Motivo Inativação': 34,
    'Dt.Desligamento': 16,
})


def _auto_map_expert(headers):
    """Mapeia colunas: tenta match exato primeiro, depois parcial."""
    mapping = {}
    hl = [(h, norm(h), norm_ascii(h)) for h in headers]
    for key, aliases in EXPERT_COL_ALIASES.items():
        # Passagem 1: match exato
        for alias in aliases:
            alias_norm = norm(alias)
            alias_ascii = norm_ascii(alias)
            for h, l, la in hl:
                if l == alias_norm or la == alias_ascii:
                    mapping[key] = h
                    break
            if key in mapping:
                break
        if key in mapping:
            continue
        # Passagem 2: match parcial (alias contido no header)
        for alias in aliases:
            if len(alias) < 5:
                continue
            alias_norm = norm(alias)
            alias_ascii = norm_ascii(alias)
            for h, l, la in hl:
                if (alias_norm in l or alias_ascii in la) and key not in mapping:
                    mapping[key] = h
                    break
            if key in mapping:
                break
    return mapping


def _prefer_expert_desc_situacao(headers):
    """No Expert, a regra de inativos deve usar a coluna descritiva."""
    preferred = {
        'desc. situacao',
        'desc situacao',
        'descricao situacao',
        'desc. situacao funcionario',
        'desc situacao funcionario',
    }
    for h in headers:
        if norm_ascii(h) in preferred:
            return h
    return None


def _to_expert_row(row, col_map):
    get = lambda k: clean_val(row.get(col_map.get(k, ''), ''))
    return {
        'Nome Unidade':               get('nome_unidade'),
        'Nome Setor':                 get('nome_setor'),
        'Nome Cargo':                 get('nome_cargo'),
        'Matrícula':                  get('matricula'),
        'Nome Funcionário':           get('nome'),
        'Dt.Nascimento':              format_date_val(row.get(col_map.get('dt_nascimento', ''), '')),
        'Sexo':                       get('sexo'),
        'Situação':                   get('situacao'),
        'Dt.Admissão':                format_date_val(row.get(col_map.get('dt_admissao', ''), '')),
        'CPF':                        format_cpf(get('cpf')),
        'CBO':                        get('cbo'),
        'Nome da Mae do Funcionário': get('nome_mae'),
        'Razão Social':               get('razao_social'),
        'CNPJ':                       format_cnpj(get('cnpj')),
        'Código Categoria (eSocial)': get('cod_categoria'),
    }


def _to_expert_inativo_row(row, col_map, motivo='Base de desligados'):
    base = _to_expert_row(row, col_map)
    situacao = base.get('Situação') or clean_val(row.get(col_map.get('situacao', ''), ''))
    return {
        'Nome Funcionário':           base.get('Nome Funcionário', ''),
        'Situação Original':          situacao,
        'Motivo Inativação':          motivo,
        'Dt.Desligamento':            format_date_val(row.get(col_map.get('dt_desligamento', ''), '')),
        'Matrícula':                  base.get('Matrícula', ''),
        'CPF':                        base.get('CPF', ''),
        'Dt.Nascimento':              base.get('Dt.Nascimento', ''),
        'Sexo':                       base.get('Sexo', ''),
        'Dt.Admissão':                base.get('Dt.Admissão', ''),
        'Nome Unidade':               base.get('Nome Unidade', ''),
        'Nome Setor':                 base.get('Nome Setor', ''),
        'Nome Cargo':                 base.get('Nome Cargo', ''),
        'CBO':                        base.get('CBO', ''),
        'Nome da Mae do Funcionário': base.get('Nome da Mae do Funcionário', ''),
        'Razão Social':               base.get('Razão Social', ''),
        'CNPJ':                       base.get('CNPJ', ''),
        'Código Categoria (eSocial)': base.get('Código Categoria (eSocial)', ''),
    }


def _write_expert_csv(rows, filepath, columns=None):
    import csv
    columns = columns or EXPERT_COLS
    with open(filepath, 'w', newline='', encoding='utf-8-sig') as f:
        writer = csv.DictWriter(f, fieldnames=columns, delimiter=';', extrasaction='ignore')
        writer.writeheader()
        for row in rows:
            writer.writerow(row)


def _expert_source_prefix(source_type):
    return 'expert_desligados_source' if source_type == 'desligados' else 'expert_ativos_source'


def _find_expert_source(session_dir, source_type='ativos'):
    prefix = _expert_source_prefix(source_type)
    return next((os.path.join(session_dir, fn)
                 for fn in os.listdir(session_dir)
                 if fn.startswith(prefix)), None)


def _read_expert_header(src, sheet_name='Planilha1'):
    ext = os.path.splitext(src)[1].lower()
    if ext == '.csv':
        for enc in ['utf-8-sig', 'latin-1', 'utf-8']:
            try:
                df = pd.read_csv(src, dtype=str, sep=None, engine='python',
                                 encoding=enc, nrows=0, na_filter=False)
                df.attrs['_encoding'] = enc
                break
            except Exception:
                df = None
        if df is None:
            raise ValueError('Erro ao ler CSV.')
    else:
        df = pd.read_excel(src, sheet_name=sheet_name, dtype=str,
                           keep_default_na=False, nrows=0)
    df.columns = [str(c).strip().lstrip('\ufeff') for c in df.columns]
    return df


def _read_expert_dataframe(src, sheet_name='Planilha1', col_map_usr=None):
    col_map_usr = col_map_usr or {}
    ext = os.path.splitext(src)[1].lower()
    hdr = _read_expert_header(src, sheet_name)
    headers = list(hdr.columns)
    auto = _auto_map_expert(headers)
    final_map = dict(auto)
    for k, v in col_map_usr.items():
        if v and v in headers:
            final_map[k] = v
    desc_situacao_col = _prefer_expert_desc_situacao(headers)
    if desc_situacao_col:
        final_map['situacao'] = desc_situacao_col

    needed = list({v for v in final_map.values() if v and v in headers})

    if ext == '.csv':
        df = None
        for enc in ['utf-8-sig', 'latin-1', 'utf-8']:
            try:
                df = pd.read_csv(src, dtype=str, sep=None, engine='python',
                                 encoding=enc, na_filter=False,
                                 usecols=needed if needed else None)
                break
            except Exception:
                continue
        if df is None:
            raise ValueError('Erro ao ler CSV.')
    else:
        df = pd.read_excel(src, sheet_name=sheet_name, dtype=str,
                           keep_default_na=False,
                           usecols=needed if needed else None)

    df.columns = [str(c).strip().lstrip('\ufeff') for c in df.columns]
    df = df[df.apply(lambda r: r.astype(str).str.strip().ne('').any(), axis=1)].reset_index(drop=True)
    for dk in ('dt_admissao', 'dt_nascimento', 'dt_desligamento'):
        col = final_map.get(dk)
        if col and col in df.columns:
            df[col] = df[col].apply(parse_mixed_dates)
    return df, final_map


def _parse_expert_filter_date(value, field_label):
    if not value:
        return None
    dt = pd.to_datetime(str(value), dayfirst=False, errors='coerce')
    if pd.isna(dt):
        raise ValueError(f'Data inválida em {field_label}.')
    return dt.normalize()


def _expert_exclusion_reason(row_d, col_map):
    if _is_pj_expert(row_d, col_map):
        return 'Pessoa Jurídica (PJ)'

    vinc_cols = [col_map.get('tipo_vinculo', ''), col_map.get('cod_categoria', '')]
    vinc = ' '.join(norm_ascii(clean_val(row_d.get(c, ''))) for c in vinc_cols if c)
    if 'contribuinte' in vinc or ('diretor' in vinc and 'fgts' in vinc):
        return 'Pessoa Jurídica (PJ)'
    if any(k in vinc for k in ['estagiario', 'estagio']):
        return 'Estagiário'
    if 'aprendiz' in vinc:
        return 'Menor Aprendiz'
    return None


def _generate_expert_ppt(ativos, inativos, excluidos, filepath):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return False
    NAVY  = RGBColor(0x0F, 0x2A, 0x56)
    TEAL  = RGBColor(0x00, 0xC2, 0xA0)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY  = RGBColor(0x64, 0x74, 0x8B)
    RED   = RGBColor(0xEF, 0x44, 0x44)
    GREEN = RGBColor(0x10, 0xB9, 0x81)
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    def rect(sl, l, t, w, h, fill=None):
        s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        s.line.fill.background()
        if fill:
            s.fill.solid()
            s.fill.fore_color.rgb = fill
        else:
            s.fill.background()
    def txt(sl, text, l, t, w, h, size=12, bold=False, color=None, align=PP_ALIGN.LEFT):
        tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        p  = tf.paragraphs[0]
        p.alignment = align
        r  = p.add_run()
        r.text = str(text)
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color or NAVY
    def metric(sl, l, t, label, value, color):
        rect(sl, l, t, 2.8, 1.3, fill=color)
        txt(sl, label, l+0.15, t+0.12, 2.5, 0.4, 9, False, WHITE)
        txt(sl, str(value), l+0.15, t+0.5, 2.5, 0.6, 30, True, WHITE)
    total   = len(ativos) + len(inativos) + len(excluidos)
    excl_e  = sum(1 for r in excluidos if r.get('_motivo') == 'Estagiário')
    excl_pj = sum(1 for r in excluidos if r.get('_motivo') == 'Pessoa Jurídica (PJ)')
    today   = __import__('datetime').date.today().strftime('%d/%m/%Y')
    sl = prs.slides.add_slide(blank)
    rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
    rect(sl, 0, 5.8, 13.33, 1.7, fill=TEAL)
    txt(sl, 'Expert', 0.6, 1.2, 12, 1.0, 52, True, WHITE)
    txt(sl, 'Relatório de Bases de Profissionais', 0.6, 2.3, 12, 0.6, 18, False, WHITE)
    txt(sl, f'Processado em {today}  ·  Stefanini Group', 0.6, 3.0, 12, 0.5, 12, False, RGBColor(0xAA,0xBB,0xCC))
    sl = prs.slides.add_slide(blank)
    rect(sl, 0, 0, 13.33, 1.0, fill=NAVY)
    txt(sl, 'Resumo Executivo', 0.4, 0.15, 10, 0.7, 20, True, WHITE)
    metric(sl, 0.4, 1.2, 'Total Analisado', total, NAVY)
    metric(sl, 3.4, 1.2, 'Ativos', len(ativos), GREEN)
    metric(sl, 6.4, 1.2, 'Inativos', len(inativos), GRAY)
    metric(sl, 9.4, 1.2, 'Removidos', len(excluidos), RED)
    prs.save(filepath)
    return True


def _generate_expert_ppt(ativos, inativos, excluidos, filepath):
    """Gera PPT executivo do cliente Expert com leitura gerencial da base."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return False

    NAVY  = RGBColor(0x0F, 0x2A, 0x56)
    TEAL  = RGBColor(0x00, 0xC2, 0xA0)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY  = RGBColor(0x64, 0x74, 0x8B)
    LGRAY = RGBColor(0xF1, 0xF5, 0xF9)
    TEXT  = RGBColor(0x1E, 0x29, 0x3B)
    RED   = RGBColor(0xEF, 0x44, 0x44)
    GREEN = RGBColor(0x10, 0xB9, 0x81)
    AMBER = RGBColor(0xF5, 0x9E, 0x0B)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def clean_text(v):
        return str(v or '').strip()

    def pct(n, d):
        return '0,0%' if not d else f'{(n / d) * 100:.1f}%'.replace('.', ',')

    def fmt(n):
        return f'{int(n):,}'.replace(',', '.')

    def count_by(rows, *keys):
        counts = {}
        for row in rows:
            value = ''
            for key in keys:
                value = clean_text(row.get(key, ''))
                if value:
                    break
            value = value or 'Nao informado'
            counts[value] = counts.get(value, 0) + 1
        return dict(sorted(counts.items(), key=lambda item: item[1], reverse=True))

    def rect(sl, l, t, w, h, fill=None, line=None):
        shp = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        if line:
            shp.line.color.rgb = line
            shp.line.width = Pt(1)
        else:
            shp.line.fill.background()
        if fill:
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
        else:
            shp.fill.background()
        return shp

    def txt(sl, text, l, t, w, h, size=16, bold=False, color=None, align=PP_ALIGN.LEFT):
        box = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or TEXT
        return box

    def title(sl, text):
        rect(sl, 0, 0, 13.33, 0.9, fill=NAVY)
        txt(sl, text, 0.45, 0.18, 12.2, 0.52, 28, True, WHITE)

    def footer(sl, page):
        txt(sl, 'BenefProcess | Expert | Stefanini Group', 0.45, 7.1, 7.2, 0.22, 8, False, GRAY)
        txt(sl, str(page), 12.35, 7.1, 0.45, 0.22, 8, False, GRAY, PP_ALIGN.RIGHT)

    def metric(sl, l, t, label, value, color, sub=''):
        rect(sl, l, t, 2.75, 1.25, fill=color)
        txt(sl, label, l + 0.15, t + 0.14, 2.45, 0.25, 11, True, WHITE)
        txt(sl, str(value), l + 0.15, t + 0.45, 2.45, 0.45, 27, True, WHITE)
        if sub:
            txt(sl, sub, l + 0.15, t + 0.95, 2.45, 0.2, 8, False, WHITE)

    def bullet(sl, text, l, t, w, size=15, color=TEXT):
        txt(sl, f'- {text}', l, t, w, 0.35, size, False, color)

    def table(sl, x, y, widths, rows, header_fill=NAVY, row_h=0.36, font=10):
        for ri, row in enumerate(rows):
            xcur = x
            fill = header_fill if ri == 0 else (LGRAY if ri % 2 else WHITE)
            color = WHITE if ri == 0 else TEXT
            bold = ri == 0
            for ci, cell in enumerate(row):
                rect(sl, xcur, y + ri * row_h, widths[ci], row_h, fill=fill, line=RGBColor(0xE2, 0xE8, 0xF0))
                align = PP_ALIGN.RIGHT if ci > 0 else PP_ALIGN.LEFT
                txt(sl, str(cell), xcur + 0.06, y + ri * row_h + 0.07, widths[ci] - 0.12, row_h - 0.08, font, bold, color, align)
                xcur += widths[ci]

    def bar(sl, label, value, total, x, y, w, color):
        txt(sl, label, x, y, 3.2, 0.22, 12, True, TEXT)
        rect(sl, x + 3.25, y + 0.02, w, 0.22, fill=RGBColor(0xE2, 0xE8, 0xF0))
        width = 0 if total == 0 else max(0.05, w * value / total)
        rect(sl, x + 3.25, y + 0.02, width, 0.22, fill=color)
        txt(sl, f'{fmt(value)} ({pct(value, total)})', x + 3.25 + w + 0.15, y - 0.02, 1.7, 0.26, 11, True, TEXT)

    total = len(ativos) + len(inativos) + len(excluidos)
    impactados = len(inativos) + len(excluidos)
    inativos_counts = count_by(inativos, 'Situação Original', 'Situacao Original', 'Situação', 'Situacao')
    excluidos_counts = count_by(excluidos, '_motivo', 'Motivo')
    empresa_counts = count_by(ativos + inativos, 'Razão Social', 'Razao Social', 'Nome Empresa')
    cargo_missing = sum(1 for r in ativos if not clean_text(r.get('Nome Cargo', '')))
    cbo_missing = sum(1 for r in ativos if not clean_text(r.get('CBO', '')))
    cpf_missing = sum(1 for r in ativos if not clean_text(r.get('CPF', '')))
    today = __import__('datetime').date.today().strftime('%d/%m/%Y')

    sl = prs.slides.add_slide(blank)
    rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
    rect(sl, 0, 5.85, 13.33, 1.65, fill=TEAL)
    txt(sl, 'Expert', 0.65, 1.15, 12, 0.85, 50, True, WHITE)
    txt(sl, 'Analise executiva da base de profissionais', 0.7, 2.1, 10.8, 0.55, 22, False, WHITE)
    txt(sl, f'Processamento local realizado em {today}', 0.7, 2.75, 7.5, 0.3, 14, False, RGBColor(0xCA, 0xDC, 0xFC))
    txt(sl, 'Objetivo: separar registros aptos para envio, controlar inativos e evidenciar exclusoes por regra de negocio.', 0.7, 6.25, 11.5, 0.55, 18, True, NAVY)

    sl = prs.slides.add_slide(blank)
    title(sl, 'A base foi saneada e separada para reduzir risco no envio')
    metric(sl, 0.55, 1.35, 'Total analisado', fmt(total), NAVY, '100% da base de entrada')
    metric(sl, 3.55, 1.35, 'Ativos para envio', fmt(len(ativos)), GREEN, pct(len(ativos), total))
    metric(sl, 6.55, 1.35, 'Inativos em controle', fmt(len(inativos)), GRAY, pct(len(inativos), total))
    metric(sl, 9.55, 1.35, 'Removidos por regra', fmt(len(excluidos)), RED, pct(len(excluidos), total))
    txt(sl, 'Leitura executiva', 0.65, 3.15, 4, 0.35, 20, True, NAVY)
    bullet(sl, f'{fmt(impactados)} registros nao seguem como ativos elegiveis e precisam ficar fora do envio principal.', 0.8, 3.75, 11.6, 16)
    bullet(sl, 'A separacao evita inclusao indevida de profissionais desligados no envio principal e preserva os afastamentos/licencas aprovados como ativos.', 0.8, 4.25, 11.6, 16)
    bullet(sl, 'Os arquivos Excel e CSV preservam o detalhe operacional; esta apresentacao resume os pontos de controle para decisao.', 0.8, 4.75, 11.6, 16)
    footer(sl, 2)

    sl = prs.slides.add_slide(blank)
    title(sl, 'A regra central usa duas bases de entrada como fonte de verdade')
    txt(sl, 'Criterio aplicado', 0.55, 1.25, 4.3, 0.35, 20, True, NAVY)
    bullet(sl, 'A base de ativos permanece como envio principal, incluindo licencas, afastamentos e demais situacoes aprovadas.', 0.75, 1.8, 5.7)
    bullet(sl, 'A base de desligados passa a alimentar o arquivo de inativos/desativados.', 0.75, 2.25, 5.7)
    bullet(sl, 'Registros da base de ativos classificados como PJ, estagiario ou aprendiz sao removidos por regra de elegibilidade.', 0.75, 2.7, 5.7)
    txt(sl, 'Entregaveis gerados', 7.0, 1.25, 4.3, 0.35, 20, True, NAVY)
    bullet(sl, 'Expert_Ativos.xlsx e Expert_Ativos.csv para envio operacional.', 7.2, 1.8, 5.4)
    bullet(sl, 'Expert_Inativos.xlsx e Expert_Inativos.csv com a situacao original preservada.', 7.2, 2.25, 5.4)
    bullet(sl, 'Apresentacao executiva com visao consolidada e pontos de atencao.', 7.2, 2.7, 5.4)
    rect(sl, 0.8, 4.1, 11.7, 1.5, fill=LGRAY)
    txt(sl, 'Ponto de controle', 1.05, 4.35, 3, 0.28, 17, True, NAVY)
    txt(sl, 'O processo separa a origem dos registros: ativos aprovados em uma base e desligados em outra, reduzindo ambiguidade na leitura de situacao.', 1.05, 4.8, 10.8, 0.45, 17, False, TEXT)
    footer(sl, 3)

    sl = prs.slides.add_slide(blank)
    title(sl, 'A maior parte da base segue apta, mas ha controles fora do envio')
    bar(sl, 'Ativos para envio', len(ativos), total, 0.75, 1.55, 6.1, GREEN)
    bar(sl, 'Inativos em controle', len(inativos), total, 0.75, 2.25, 6.1, GRAY)
    bar(sl, 'Removidos por regra', len(excluidos), total, 0.75, 2.95, 6.1, RED)
    txt(sl, 'O que isso significa', 0.9, 4.15, 3.6, 0.28, 19, True, NAVY)
    bullet(sl, 'A base final de ativos representa o universo elegivel apos saneamento.', 1.05, 4.65, 5.6)
    bullet(sl, 'Inativos nao sao erro de cadastro: sao registros que requerem controle separado.', 1.05, 5.1, 5.6)
    bullet(sl, 'Removidos representam regras objetivas que impedem envio ao cliente.', 1.05, 5.55, 5.6)
    txt(sl, 'Distribuicao por empresa', 7.3, 1.25, 4.3, 0.35, 20, True, NAVY)
    emp_rows = [['Empresa', 'Qtd.']] + [[k[:32], fmt(v)] for k, v in list(empresa_counts.items())[:7]]
    table(sl, 7.3, 1.75, [3.55, 1.2], emp_rows, font=9)
    footer(sl, 4)

    sl = prs.slides.add_slide(blank)
    title(sl, 'Os inativos vêm da base de desligados enviada separadamente')
    rows = [['Situacao original', 'Qtd.', '% base', '% inativos']]
    for k, v in list(inativos_counts.items())[:9]:
        rows.append([k, fmt(v), pct(v, total), pct(v, len(inativos))])
    table(sl, 0.75, 1.25, [4.5, 1.0, 1.35, 1.55], rows, font=9)
    txt(sl, 'Interpretacao', 9.1, 1.3, 3, 0.28, 20, True, NAVY)
    bullet(sl, 'Afastamentos, licencas e situacoes aprovadas permanecem na base de ativos.', 9.15, 1.85, 3.6, 14)
    bullet(sl, 'A preservacao da situacao original dos desligados reduz retrabalho na validacao com RH.', 9.15, 2.45, 3.6, 14)
    bullet(sl, 'O detalhe nominativo permanece no arquivo Expert_Inativos.', 9.15, 3.05, 3.6, 14)
    footer(sl, 5)

    sl = prs.slides.add_slide(blank)
    title(sl, 'As exclusoes removem perfis fora do escopo do envio')
    rows = [['Motivo da remocao', 'Qtd.', '% base', 'Tratamento']]
    treatment = {
        'Pessoa Jurídica (PJ)': 'Nao elegivel',
        'Pessoa Juridica (PJ)': 'Nao elegivel',
        'Estagiário': 'Fora da regra',
        'Estagiario': 'Fora da regra',
        'Menor Aprendiz': 'Fora da regra',
    }
    for k, v in list(excluidos_counts.items())[:8]:
        rows.append([k, fmt(v), pct(v, total), treatment.get(k, 'Revisar regra')])
    table(sl, 0.8, 1.35, [3.4, 1.0, 1.35, 2.0], rows, font=10)
    txt(sl, 'Impacto para a diretoria', 8.5, 1.35, 3.7, 0.28, 20, True, NAVY)
    bullet(sl, 'A remocao protege o envio contra grupos nao previstos no contrato ou na politica de beneficios.', 8.6, 1.9, 3.7, 14)
    bullet(sl, 'A regra e objetiva e auditavel por motivo, sem exclusao manual pontual.', 8.6, 2.65, 3.7, 14)
    bullet(sl, 'Mudancas de politica podem ser refletidas no processamento futuro sem alterar os arquivos historicos.', 8.6, 3.35, 3.7, 14)
    footer(sl, 6)

    sl = prs.slides.add_slide(blank)
    title(sl, 'Controles de qualidade deixam a base pronta para uso operacional')
    metric(sl, 0.7, 1.35, 'Ativos sem CPF', fmt(cpf_missing), GREEN if cpf_missing == 0 else AMBER, 'controle cadastral')
    metric(sl, 3.7, 1.35, 'Ativos sem cargo', fmt(cargo_missing), GREEN if cargo_missing == 0 else AMBER, 'qualidade de perfil')
    metric(sl, 6.7, 1.35, 'Ativos sem CBO', fmt(cbo_missing), GREEN if cbo_missing == 0 else AMBER, 'classificacao ocupacional')
    txt(sl, 'Controles aplicados no processamento', 0.8, 3.25, 5.2, 0.3, 20, True, NAVY)
    bullet(sl, 'Leitura estruturada da base e selecao das colunas necessarias.', 1.0, 3.8, 5.8)
    bullet(sl, 'Separacao entre ativos, inativos e removidos antes da geracao dos arquivos.', 1.0, 4.25, 5.8)
    bullet(sl, 'Geracao de Excel e CSV para conferencia e envio.', 1.0, 4.7, 5.8)
    txt(sl, 'Governanca recomendada', 7.2, 3.25, 4.2, 0.3, 20, True, NAVY)
    bullet(sl, 'Guardar os arquivos gerados como evidencia da competencia processada.', 7.4, 3.8, 5.0)
    bullet(sl, 'Validar mudancas de regra com Gente & Cultura antes de novo ciclo.', 7.4, 4.25, 5.0)
    bullet(sl, 'Tratar inativos como controle separado, nao como removidos.', 7.4, 4.7, 5.0)
    footer(sl, 7)

    sl = prs.slides.add_slide(blank)
    title(sl, 'Decisao recomendada e proximos passos')
    txt(sl, 'Recomendacao', 0.75, 1.25, 4, 0.3, 22, True, NAVY)
    txt(sl, 'Aprovar o envio da base de ativos e manter inativos e removidos segregados como evidencia de controle.', 0.8, 1.75, 11.4, 0.55, 22, True, TEXT)
    txt(sl, 'Proximos passos', 0.75, 3.0, 4, 0.3, 22, True, NAVY)
    bullet(sl, 'Enviar Expert_Ativos para o fluxo operacional aprovado.', 1.0, 3.55, 10.8, 17)
    bullet(sl, 'Compartilhar Expert_Inativos com RH para acompanhamento de retorno, licenca, afastamento ou alteracao cadastral.', 1.0, 4.05, 10.8, 17)
    bullet(sl, 'Revisar a lista de removidos apenas se houver mudanca formal na regra de elegibilidade.', 1.0, 4.55, 10.8, 17)
    bullet(sl, 'Arquivar a apresentacao e os arquivos de saida como trilha de auditoria do ciclo.', 1.0, 5.05, 10.8, 17)
    rect(sl, 0.85, 6.0, 11.65, 0.65, fill=LGRAY)
    txt(sl, 'Mensagem-chave: o processo entrega uma base utilizavel, auditavel e separada por finalidade, reduzindo risco de envio indevido.', 1.05, 6.18, 11.1, 0.3, 16, True, NAVY)
    footer(sl, 8)

    prs.save(filepath)
    return True


@app.route('/api/upload-expert', methods=['POST'])
def api_upload_expert():
    if 'file' not in request.files:
        return jsonify({'error': 'Nenhum arquivo enviado.'}), 400
    f = request.files['file']
    if not f.filename.lower().endswith(('.xlsx', '.xls', '.csv')):
        return jsonify({'error': 'Formato inválido. Envie .xlsx, .xls ou .csv.'}), 400
    source_type = request.form.get('source_type', 'ativos')
    if source_type not in ('ativos', 'desligados'):
        source_type = 'ativos'
    sid = request.form.get('session_id') or str(uuid.uuid4())
    session_dir = os.path.join(TEMP_DIR, sid)
    os.makedirs(session_dir, exist_ok=True)
    ext = os.path.splitext(f.filename)[1].lower()
    upload_path = os.path.join(session_dir, f'{_expert_source_prefix(source_type)}{ext}')
    f.save(upload_path)
    try:
        if ext == '.csv':
            sheet_names = ['Planilha1']
        else:
            sheet_names = pd.ExcelFile(upload_path).sheet_names
        return jsonify({'session_id': sid, 'sheet_names': sheet_names,
                        'filename': f.filename, 'ext': ext})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/get-columns-expert', methods=['POST'])
def api_get_columns_expert():
    data        = request.json
    sid         = data.get('session_id')
    sheet_name  = data.get('sheet_name', 'Planilha1')
    source_type = data.get('source_type', 'ativos')
    session_dir = os.path.join(TEMP_DIR, sid)
    if not os.path.isdir(session_dir):
        return jsonify({'error': 'Sessão expirada.'}), 400
    src = _find_expert_source(session_dir, source_type)
    if not src:
        return jsonify({'error': 'Arquivo não encontrado.'}), 400
    try:
        df = _read_expert_header(src, sheet_name)
        cols = [str(c).strip() for c in df.columns]
        return jsonify({'headers': cols, 'auto_map': _auto_map_expert(cols)})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/process-expert', methods=['POST'])
def api_process_expert():
    data        = request.json
    sid         = data.get('session_id')
    active_sheet_name   = data.get('active_sheet_name') or data.get('sheet_name', 'Planilha1')
    inactive_sheet_name = data.get('inactive_sheet_name') or data.get('sheet_name', 'Planilha1')
    col_map_usr = data.get('col_map', {})
    inactive_start_date = data.get('inactive_start_date') or ''
    inactive_end_date   = data.get('inactive_end_date') or ''
    session_dir = os.path.join(TEMP_DIR, sid)
    if not os.path.isdir(session_dir):
        return jsonify({'error': 'Sessão expirada.'}), 400
    active_src = _find_expert_source(session_dir, 'ativos')
    inactive_src = _find_expert_source(session_dir, 'desligados')
    if not active_src:
        return jsonify({'error': 'Envie a base de ativos do Expert antes de processar.'}), 400
    if not inactive_src:
        return jsonify({'error': 'Envie a base de desligados do Expert antes de processar.'}), 400
    try:
        active_df, active_map = _read_expert_dataframe(active_src, active_sheet_name, col_map_usr)
        inactive_df, inactive_map = _read_expert_dataframe(inactive_src, inactive_sheet_name, col_map_usr)
        ativos_rows, inativos_rows, excluidos_rows = [], [], []

        if inactive_start_date or inactive_end_date:
            deslig_col = inactive_map.get('dt_desligamento', '')
            if not deslig_col or deslig_col not in inactive_df.columns:
                return jsonify({
                    'error': 'Para filtrar desligados por período, mapeie a coluna "Data de desligamento/demissão".'
                }), 400
            start_dt = _parse_expert_filter_date(inactive_start_date, 'Data inicial dos desligados')
            end_dt = _parse_expert_filter_date(inactive_end_date, 'Data final dos desligados')
            if start_dt is not None and end_dt is not None and start_dt > end_dt:
                return jsonify({'error': 'A data inicial dos desligados não pode ser maior que a data final.'}), 400

            deslig_dates = pd.to_datetime(inactive_df[deslig_col], errors='coerce').dt.normalize()
            mask = deslig_dates.notna()
            if start_dt is not None:
                mask &= deslig_dates >= start_dt
            if end_dt is not None:
                mask &= deslig_dates <= end_dt
            inactive_df = inactive_df[mask].reset_index(drop=True)

        # A base de ativos nao separa mais inativos por Desc. Situacao:
        # todas as situacoes aprovadas pelo cliente seguem como ativas, exceto
        # perfis removidos por regra de elegibilidade.
        for row in active_df.to_dict('records'):
            motivo = _expert_exclusion_reason(row, active_map)
            if motivo:
                row['_motivo'] = motivo
                excluidos_rows.append(row)
            else:
                ativos_rows.append(_to_expert_row(row, active_map))

        # A base de desligados e a fonte de verdade dos inativos/desativados.
        inativos_rows = [
            _to_expert_inativo_row(r, inactive_map, 'Base de desligados')
            for r in inactive_df.to_dict('records')
        ]

        write_fast_excel(ativos_rows,   os.path.join(session_dir, 'expert_ativos.xlsx'),   'Ativos',   EXPERT_COLS)
        write_fast_excel(inativos_rows, os.path.join(session_dir, 'expert_inativos.xlsx'), 'Inativos', EXPERT_INATIVOS_COLS)
        _write_expert_csv(ativos_rows,    os.path.join(session_dir, 'expert_ativos.csv'))
        _write_expert_csv(inativos_rows,  os.path.join(session_dir, 'expert_inativos.csv'), EXPERT_INATIVOS_COLS)

        ppt_ok = _generate_expert_ppt(
            ativos_rows,
            inativos_rows,
            excluidos_rows,
            os.path.join(session_dir, 'expert_ppt.pptx')
        )

        def count_reasons(rows):
            c = {}
            for r in rows:
                m = r.get('_motivo', 'Outro')
                c[m] = c.get(m, 0) + 1
            return c

        return jsonify({
            'total': len(active_df) + len(inactive_df),
            'ativos': len(ativos_rows),
            'inativos': len(inativos_rows),
            'excluidos': len(excluidos_rows),
            'reasons': count_reasons(excluidos_rows),
            'ppt_available': ppt_ok,
        })

        ext = os.path.splitext(src)[1].lower()

        # 1. Lê só os cabeçalhos para montar o mapeamento
        if ext == '.csv':
            _hdr = None
            _enc = 'latin-1'
            for enc in ['utf-8-sig', 'latin-1', 'utf-8']:
                try:
                    _hdr = pd.read_csv(src, dtype=str, sep=None, engine='python',
                                       encoding=enc, nrows=0, na_filter=False)
                    _enc = enc
                    break
                except Exception:
                    continue
            if _hdr is None:
                return jsonify({'error': 'Erro ao ler CSV.'}), 500
        else:
            _hdr = pd.read_excel(src, sheet_name=sheet_name, dtype=str,
                                 keep_default_na=False, nrows=0)

        _hdr.columns = [str(c).strip().lstrip('﻿') for c in _hdr.columns]
        auto      = _auto_map_expert(list(_hdr.columns))
        final_map = {**auto, **{k: v for k, v in col_map_usr.items() if v}}
        desc_situacao_col = _prefer_expert_desc_situacao(list(_hdr.columns))
        if desc_situacao_col:
            final_map['situacao'] = desc_situacao_col

        # 2. Lê apenas as colunas necessárias (3x mais rápido)
        _needed = list({v for v in final_map.values() if v and v in _hdr.columns})

        if ext == '.csv':
            df = None
            for enc in ['utf-8-sig', 'latin-1', 'utf-8']:
                try:
                    df = pd.read_csv(src, dtype=str, sep=None, engine='python',
                                     encoding=enc, na_filter=False,
                                     usecols=_needed if _needed else None)
                    break
                except Exception:
                    continue
            if df is None:
                return jsonify({'error': 'Erro ao ler CSV.'}), 500
        else:
            df = pd.read_excel(src, sheet_name=sheet_name, dtype=str,
                               keep_default_na=False,
                               usecols=_needed if _needed else None)
        df.columns = [str(c).strip().lstrip('﻿') for c in df.columns]
        df = df[df.apply(lambda r: r.astype(str).str.strip().ne('').any(), axis=1)].reset_index(drop=True)
        for dk in ('dt_admissao', 'dt_nascimento'):
            col = final_map.get(dk)
            if col and col in df.columns:
                df[col] = df[col].apply(parse_mixed_dates)
        situacao_col = final_map.get('situacao', '')
        cod_cat_col  = final_map.get('cod_categoria', '')
        ativos_rows, inativos_rows, excluidos_rows = [], [], []
        # ── Processamento vetorizado (100x mais rápido para 95k+ linhas) ──
        sit_col  = situacao_col if situacao_col and situacao_col in df.columns else None
        vinc_col = cod_cat_col  if cod_cat_col  and cod_cat_col  in df.columns else None

        if not sit_col:
            return jsonify({
                'error': 'Mapeie a coluna "Desc. Situação" para separar ativos e inativos do Expert.'
            }), 400

        sit  = df[sit_col].map(norm_ascii).fillna('')  if sit_col  else pd.Series([EXPERT_SITUACAO_ATIVA]*len(df))
        vinc = df[vinc_col].map(norm_ascii).fillna('') if vinc_col else pd.Series(['']*len(df))

        inativo_mask = sit != norm_ascii(EXPERT_SITUACAO_ATIVA)
        active_mask  = ~inativo_mask

        pj_mask  = active_mask & (
            vinc.str.contains('contribuinte', na=False) |
            (vinc.str.contains('diretor', na=False) & vinc.str.contains('fgts', na=False))
        )
        est_mask = active_mask & ~pj_mask & (
            vinc.str.contains('estagiario|estagiário|estágio|estagio', na=False)
        )
        apr_mask = active_mask & ~pj_mask & ~est_mask & (
            vinc.str.contains('aprendiz', na=False)
        )
        ativ_mask = active_mask & ~pj_mask & ~est_mask & ~apr_mask

        def _tag(sub_df, motivo):
            recs = sub_df.to_dict('records')
            for r in recs: r['_motivo'] = motivo
            return recs

        ativos_rows    = [_to_expert_row(r, final_map) for r in df[ativ_mask].to_dict('records')]
        inativos_rows  = [_to_expert_inativo_row(r, final_map) for r in df[inativo_mask].to_dict('records')]
        excluidos_rows = (_tag(df[pj_mask],  'Pessoa Jurídica (PJ)') +
                          _tag(df[est_mask],  'Estagiário') +
                          _tag(df[apr_mask],  'Menor Aprendiz'))

        write_fast_excel(ativos_rows,   os.path.join(session_dir, 'expert_ativos.xlsx'),   'Ativos',   EXPERT_COLS)
        write_fast_excel(inativos_rows, os.path.join(session_dir, 'expert_inativos.xlsx'), 'Inativos', EXPERT_INATIVOS_COLS)
        _write_expert_csv(ativos_rows,    os.path.join(session_dir, 'expert_ativos.csv'))
        _write_expert_csv(inativos_rows,  os.path.join(session_dir, 'expert_inativos.csv'), EXPERT_INATIVOS_COLS)

        ppt_ok = _generate_expert_ppt(
            ativos_rows,
            inativos_rows,
            excluidos_rows,
            os.path.join(session_dir, 'expert_ppt.pptx')
        )

        def count_reasons(rows):
            c = {}
            for r in rows:
                m = r.get('_motivo', 'Outro')
                c[m] = c.get(m, 0) + 1
            return c
        return jsonify({
            'total': len(df), 'ativos': len(ativos_rows),
            'inativos': len(inativos_rows), 'excluidos': len(excluidos_rows),
            'reasons': count_reasons(excluidos_rows),
            'ppt_available': ppt_ok,
        })
    except Exception as e:
        import traceback
        return jsonify({'error': str(e), 'detail': traceback.format_exc()}), 500


@app.route('/api/download/<session_id>/expert_ativos_excel')
def download_expert_ativos_excel(session_id):
    p = os.path.join(TEMP_DIR, session_id, 'expert_ativos.xlsx')
    if not os.path.exists(p): return 'Não encontrado.', 404
    return send_file(p, as_attachment=True, download_name='Expert_Ativos.xlsx')

@app.route('/api/download/<session_id>/expert_inativos_excel')
def download_expert_inativos_excel(session_id):
    p = os.path.join(TEMP_DIR, session_id, 'expert_inativos.xlsx')
    if not os.path.exists(p): return 'Não encontrado.', 404
    return send_file(p, as_attachment=True, download_name='Expert_Inativos.xlsx')

@app.route('/api/download/<session_id>/expert_ativos_csv')
def download_expert_ativos_csv(session_id):
    p = os.path.join(TEMP_DIR, session_id, 'expert_ativos.csv')
    if not os.path.exists(p): return 'Não encontrado.', 404
    return send_file(p, as_attachment=True, download_name='Expert_Ativos.csv')

@app.route('/api/download/<session_id>/expert_inativos_csv')
def download_expert_inativos_csv(session_id):
    p = os.path.join(TEMP_DIR, session_id, 'expert_inativos.csv')
    if not os.path.exists(p): return 'Não encontrado.', 404
    return send_file(p, as_attachment=True, download_name='Expert_Inativos.csv')

@app.route('/api/download/<session_id>/expert_ppt')
def download_expert_ppt(session_id):
    session_dir = os.path.join(TEMP_DIR, session_id)
    p = os.path.join(session_dir, 'expert_ppt.pptx')
    if not os.path.exists(p):
        ativos_path = os.path.join(session_dir, 'expert_ativos.xlsx')
        inativos_path = os.path.join(session_dir, 'expert_inativos.xlsx')
        if not os.path.exists(ativos_path) or not os.path.exists(inativos_path):
            return 'PPT não encontrado. Processe a base Expert antes de baixar a apresentação.', 404
        try:
            ativos = pd.read_excel(ativos_path, dtype=str).fillna('').to_dict('records')
            inativos = pd.read_excel(inativos_path, dtype=str).fillna('').to_dict('records')
            if not _generate_expert_ppt(ativos, inativos, [], p):
                return 'Dependência ausente para gerar PPT. Instale python-pptx.', 500
        except Exception as e:
            return f'Erro ao gerar PPT Expert: {str(e)}', 500
    return send_file(p, as_attachment=True, download_name='Expert_Apresentacao_Executiva.pptx')

# ══════════════════════════════════════════════════════════════════════════════
if __name__ == '__main__':
    print('=' * 50)
    print('  BenefProcess — Tratamento de Benefícios')
    print('  Stefanini · Gente e Cultura')
    print('=' * 50)
    print('  Acesse no navegador: http://localhost:5050')
    print('=' * 50)
    app.run(debug=False, port=5050)
