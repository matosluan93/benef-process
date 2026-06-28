"""
ppt_generator.py
Geração do PPT de auditoria de exclusões de benefícios.
Lê o auditoria.xlsx já gerado pelo BenefProcess e produz um .pptx.
"""
import re
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from lxml import etree


# ── Cores ─────────────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x0F, 0x2A, 0x56)
BLUE  = RGBColor(0x1A, 0x4A, 0x8A)
TEAL  = RGBColor(0x08, 0x91, 0xB2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0x64, 0x74, 0x8B)
LGRAY = RGBColor(0xF1, 0xF5, 0xF9)
RED   = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xD9, 0x77, 0x06)
GREEN = RGBColor(0x05, 0x96, 0x69)
TEXT  = RGBColor(0x1E, 0x29, 0x3B)
LIGHT = RGBColor(0xE8, 0xF0, 0xFA)

W = Inches(10)
H = Inches(5.625)


# ── Helpers ───────────────────────────────────────────────────────────────────
def rgb_hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None, radius=0):
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1 if radius == 0 else 5,  # 1=rect, 5=rounded rect
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if radius and hasattr(shape, "adjustments"):
        try:
            shape.adjustments[0] = 0.08
        except Exception:
            pass
    return shape


def add_text(slide, text, x, y, w, h, size=12, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txBox


def set_chart_datalabel_size(chart, size_pt=11, bold=True):
    """Aplica tamanho de fonte nos data labels via XML."""
    try:
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
        for dLbls in chart._element.iter(f'{{{ns}}}dLbls'):
            txPr = etree.SubElement(dLbls, f'{{{ns}}}txPr')
            bodyPr = etree.SubElement(txPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
            lstStyle = etree.SubElement(txPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle')
            p = etree.SubElement(txPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}p')
            pPr = etree.SubElement(p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
            defRPr = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr')
            defRPr.set('sz', str(size_pt * 100))
            if bold:
                defRPr.set('b', '1')
    except Exception:
        pass


def set_series_color(series, rgb: RGBColor):
    """Aplica cor a uma série de gráfico via XML."""
    try:
        ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        spPr = etree.SubElement(series._element, f'{{{ns_c}}}spPr')
        solidFill = etree.SubElement(
            etree.SubElement(spPr, f'{{{ns_a}}}solidFill'),
            f'{{{ns_a}}}srgbClr'
        )
        solidFill.set('val', rgb_hex(rgb))
    except Exception:
        pass


def show_datalabels(chart, show_val=True):
    """Ativa data labels num gráfico via XML."""
    try:
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
        for plotArea in chart._element.iter(f'{{{ns}}}plotArea'):
            for ser in plotArea.iter(f'{{{ns}}}ser'):
                dLbls = etree.SubElement(ser, f'{{{ns}}}dLbls')
                showVal = etree.SubElement(dLbls, f'{{{ns}}}showVal')
                showVal.set('val', '1' if show_val else '0')
                showLegKey = etree.SubElement(dLbls, f'{{{ns}}}showLegKey')
                showLegKey.set('val', '0')
                showCatName = etree.SubElement(dLbls, f'{{{ns}}}showCatName')
                showCatName.set('val', '0')
                showSerName = etree.SubElement(dLbls, f'{{{ns}}}showSerName')
                showSerName.set('val', '0')
                showPercent = etree.SubElement(dLbls, f'{{{ns}}}showPercent')
                showPercent.set('val', '0')
                showBubbleSize = etree.SubElement(dLbls, f'{{{ns}}}showBubbleSize')
                showBubbleSize.set('val', '0')
    except Exception:
        pass


# ── Leitura dos dados ─────────────────────────────────────────────────────────
def load_stats(auditoria_path):
    df = pd.read_excel(auditoria_path, dtype=str)

    total_excl   = len(df)
    por_processo = df['Processo'].value_counts().to_dict()
    por_motivo   = df['Motivo da exclusão'].value_counts().to_dict()
    por_empresa  = df['Nome da Empresa'].value_counts().head(6).to_dict()

    cross = pd.crosstab(df['Motivo da exclusão'], df['Processo']).to_dict()

    return {
        'total_excluidos': total_excl,
        'por_processo': por_processo,
        'por_motivo': por_motivo,
        'por_empresa': por_empresa,
        'cross': cross,
    }


# ── Slide 1: Capa ─────────────────────────────────────────────────────────────
def slide_capa(prs, data):
    layout = prs.slide_layouts[6]  # blank
    s = prs.slides.add_slide(layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY

    add_rect(s, 0.6, 1.0, 2.2, 0.38, TEAL)
    add_text(s, "GENTE & CULTURA", 0.6, 1.0, 2.2, 0.38, size=9, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, "Auditoria de Exclusões", 0.6, 1.55, 8.5, 1.0,
             size=40, bold=True, color=WHITE, font="Cambria")
    add_text(s, "Processamento de Benefícios — Junho 2026", 0.6, 2.55, 8.5, 0.55,
             size=18, color=RGBColor(0xCA, 0xDC, 0xFC))
    add_text(s, "TotalPass  ·  Wellhub  ·  New Value", 0.6, 3.2, 6, 0.4,
             size=13, color=RGBColor(0x94, 0xA3, 0xB8))
    add_text(s, "BenefProcess  ·  Stefanini Group", 0.6, 4.95, 8, 0.35,
             size=10, color=RGBColor(0x47, 0x56, 0x69))


# ── Slide 2: Visão Geral ──────────────────────────────────────────────────────
def slide_overview(prs, data, total_base=18218):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE

    add_text(s, "Visão Geral do Processamento", 0.5, 0.28, 9, 0.65,
             size=28, bold=True, color=NAVY, font="Cambria")

    excl  = data['total_excluidos']
    taxa  = excl / total_base * 100

    cards = [
        (f"{total_base:,}".replace(',', '.'), "Registros na Base",  "colaboradores processados", NAVY),
        (f"{excl:,}".replace(',', '.'),       "Excluídos",           "não atenderam os critérios", RED),
        (f"{taxa:.1f}%".replace('.', ','),    "Taxa de Exclusão",    "do total da base",           AMBER),
    ]
    for i, (val, lbl, sub, cor) in enumerate(cards):
        x = 0.5 + i * 3.1
        add_rect(s, x, 1.05, 2.8, 1.9, cor)
        add_text(s, val, x, 1.15, 2.8, 0.9, size=36, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font="Cambria")
        add_text(s, lbl, x, 2.07, 2.8, 0.35, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(s, sub, x, 2.44, 2.8, 0.3, size=9,
                 color=RGBColor(0xCA, 0xDC, 0xFC), align=PP_ALIGN.CENTER)

    # Tabela por processo
    tp = data['por_processo']
    add_text(s, "Exclusões por Processo", 0.5, 3.15, 9, 0.38,
             size=14, bold=True, color=NAVY, font="Cambria")

    procs = [
        ("TotalPass",  tp.get('TotalPass', 0),  NAVY),
        ("Wellhub",    tp.get('Wellhub', 0),    TEAL),
        ("New Value",  tp.get('New Value', 0),  AMBER),
    ]
    for i, (nome, qtd, cor) in enumerate(procs):
        x = 0.5 + i * 3.1
        add_rect(s, x, 3.6, 2.8, 1.7, LGRAY)
        add_text(s, nome, x, 3.7, 2.8, 0.38, size=13, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
        add_text(s, f"{qtd:,}".replace(',', '.'), x, 4.1, 2.8, 0.55,
                 size=28, bold=True, color=cor, align=PP_ALIGN.CENTER, font="Cambria")
        pct = qtd / excl * 100 if excl else 0
        add_text(s, f"{pct:.1f}% das exclusões".replace('.', ','), x, 4.68, 2.8, 0.3,
                 size=10, color=GRAY, align=PP_ALIGN.CENTER)


# ── Slide 3: Motivos ──────────────────────────────────────────────────────────
def slide_motivos(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE

    add_text(s, "Motivos de Exclusão — Ranking Completo", 0.5, 0.28, 9, 0.65,
             size=24, bold=True, color=NAVY, font="Cambria")
    add_text(s, f"Distribuição dos {data['total_excluidos']:,} registros excluídos por categoria".replace(',', '.'),
             0.5, 0.88, 9, 0.3, size=12, color=GRAY)

    motivos_cores = {
        "E-mail sem domínio Stefanini": RGBColor(0xE5, 0x3E, 0x3E),
        "Pessoa Jurídica (PJ)":         RGBColor(0xD9, 0x77, 0x06),
        "E-mail inválido ou ausente":   RGBColor(0x7C, 0x3A, 0xED),
        "E-mail sem domínio IHM":       TEAL,
        "Estagiário":                   GREEN,
        "Aprendiz":                     NAVY,
        "Data de admissão futura":      GRAY,
    }
    pm = data['por_motivo']
    excl = data['total_excluidos']
    labels = list(motivos_cores.keys())
    values = [pm.get(k, 0) for k in labels]
    cores  = list(motivos_cores.values())

    # Gráfico de barras horizontal
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series('Exclusões', values)
    chart = s.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.4), Inches(1.2), Inches(6.2), Inches(4.1),
        cd
    ).chart

    # Cores por série via XML
    try:
        ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ser_el = list(chart._element.iter(f'{{{ns_c}}}ser'))[0]
        for i, cor in enumerate(reversed(cores)):
            dp = etree.SubElement(ser_el, f'{{{ns_c}}}dPt')
            idx = etree.SubElement(dp, f'{{{ns_c}}}idx')
            idx.set('val', str(i))
            spPr = etree.SubElement(dp, f'{{{ns_c}}}spPr')
            sf = etree.SubElement(
                etree.SubElement(spPr, f'{{{ns_a}}}solidFill'),
                f'{{{ns_a}}}srgbClr')
            sf.set('val', rgb_hex(cor))
    except Exception:
        pass

    show_datalabels(chart)
    set_chart_datalabel_size(chart, 11, True)
    chart.has_legend = False

    # Tabela lateral
    add_text(s, "Motivo", 6.8, 1.2, 1.6, 0.35, size=10, bold=True,
             color=WHITE, font="Calibri")
    add_rect(s, 6.8, 1.2, 2.8, 0.35, NAVY)
    add_text(s, "Motivo", 6.8, 1.2, 1.6, 0.35, size=10, bold=True, color=WHITE)
    add_text(s, "Qtd", 8.4, 1.2, 0.7, 0.35, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "%", 9.1, 1.2, 0.5, 0.35, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, (k, cor) in enumerate(motivos_cores.items()):
        y = 1.55 + i * 0.52
        bg = LGRAY if i % 2 == 0 else WHITE
        add_rect(s, 6.8, y, 2.8, 0.5, bg)
        qtd = pm.get(k, 0)
        pct = qtd / excl * 100 if excl else 0
        lbl = k.replace("E-mail sem domínio ", "s/ ").replace("E-mail inválido ou ausente", "E-mail inválido")
        add_text(s, lbl, 6.82, y + 0.05, 1.55, 0.42, size=9, color=TEXT)
        add_text(s, f"{qtd:,}".replace(',', '.'), 8.4, y + 0.05, 0.7, 0.42,
                 size=11, bold=True, color=cor, align=PP_ALIGN.CENTER)
        add_text(s, f"{pct:.1f}%".replace('.', ','), 9.1, y + 0.05, 0.5, 0.42,
                 size=9, color=GRAY, align=PP_ALIGN.CENTER)


# ── Slide 4: E-mail ───────────────────────────────────────────────────────────
def slide_email(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)

    pm = data['por_motivo']
    excl = data['total_excluidos']
    n_stef = pm.get("E-mail sem domínio Stefanini", 0)
    n_ihm  = pm.get("E-mail sem domínio IHM", 0)
    n_inv  = pm.get("E-mail inválido ou ausente", 0)
    total_email = n_stef + n_ihm + n_inv
    pct_email = total_email / excl * 100 if excl else 0

    add_text(s, "Análise: Problemas de E-mail", 0.5, 0.28, 9, 0.65,
             size=28, bold=True, color=NAVY, font="Cambria")
    add_text(s, f"{total_email:,} exclusões por e-mail — {pct_email:.1f}% do total".replace(',', '.').replace('.', ',', 1),
             0.5, 0.88, 9, 0.3, size=12, color=GRAY)

    cats = [
        ("E-mail sem domínio Stefanini", n_stef, RED,
         "E-mails pessoais ou externos em vez de @stefanini.com"),
        ("E-mail sem domínio IHM",       n_ihm,  TEAL,
         "Gmail / hotmail em vez de @ihm.com.br"),
        ("E-mail inválido ou ausente",    n_inv,  RGBColor(0x7C, 0x3A, 0xED),
         "Campo vazio ou formato inválido"),
    ]

    # Gráfico de barras horizontal
    cd = CategoryChartData()
    cd.categories = [c[0].replace("E-mail sem domínio ", "s/ ") for c in cats]
    cd.add_series('Exclusões', [c[1] for c in cats])
    chart = s.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.4), Inches(1.2), Inches(4.8), Inches(2.8),
        cd
    ).chart
    try:
        ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ser_el = list(chart._element.iter(f'{{{ns_c}}}ser'))[0]
        for i, (_, _, cor, _) in enumerate(cats):
            dp = etree.SubElement(ser_el, f'{{{ns_c}}}dPt')
            idx = etree.SubElement(dp, f'{{{ns_c}}}idx')
            idx.set('val', str(i))
            spPr = etree.SubElement(dp, f'{{{ns_c}}}spPr')
            sf = etree.SubElement(
                etree.SubElement(spPr, f'{{{ns_a}}}solidFill'),
                f'{{{ns_a}}}srgbClr')
            sf.set('val', rgb_hex(cor))
    except Exception:
        pass
    show_datalabels(chart)
    set_chart_datalabel_size(chart, 13, True)
    chart.has_legend = False

    # Cards
    for i, (titulo, qtd, cor, desc) in enumerate(cats):
        y = 1.2 + i * 1.12
        add_rect(s, 5.4, y, 4.1, 0.96, WHITE)
        add_rect(s, 5.4, y, 0.07, 0.96, cor)
        pct = qtd / excl * 100 if excl else 0
        add_text(s, titulo, 5.55, y + 0.06, 2.6, 0.3, size=11, bold=True, color=NAVY)
        add_text(s, desc, 5.55, y + 0.42, 2.6, 0.44, size=9, color=GRAY)
        add_text(s, f"{qtd:,}".replace(',', '.'), 8.2, y + 0.06, 1.2, 0.44,
                 size=24, bold=True, color=cor, align=PP_ALIGN.RIGHT, font="Cambria")
        add_text(s, f"{pct:.1f}% do total".replace('.', ','), 8.2, y + 0.58, 1.2, 0.28,
                 size=9, color=GRAY, align=PP_ALIGN.RIGHT)

    add_rect(s, 0.4, 4.3, 9.2, 0.42, RGBColor(0xFE, 0xF3, 0xC7))
    add_text(s, f"Total de exclusões por e-mail: {total_email:,}  ({pct_email:.1f}% de todas as exclusões)".replace(',', '.').replace('.', ',', 1),
             0.55, 4.34, 8.8, 0.34, size=11, bold=True, color=RGBColor(0x92, 0x40, 0x0E))


# ── Slide 5: Vínculo ──────────────────────────────────────────────────────────
def slide_vinculo(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE

    pm = data['por_motivo']
    excl = data['total_excluidos']
    n_pj  = pm.get("Pessoa Jurídica (PJ)", 0)
    n_est = pm.get("Estagiário", 0)
    n_apr = pm.get("Aprendiz", 0)
    total_vinc = n_pj + n_est + n_apr

    add_text(s, "Análise: Tipo de Vínculo Inelegível", 0.5, 0.28, 9, 0.65,
             size=24, bold=True, color=NAVY, font="Cambria")
    add_text(s, f"{total_vinc:,} exclusões por vínculo inelegível — {total_vinc/excl*100:.1f}% do total".replace(',', '.').replace('.', ',', 1),
             0.5, 0.88, 9, 0.3, size=12, color=GRAY)

    cd = CategoryChartData()
    cd.categories = ["Pessoa Jurídica (PJ)", "Estagiário", "Aprendiz"]
    cd.add_series("Vínculo", [n_pj, n_est, n_apr])

    chart = s.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(0.4), Inches(1.2), Inches(4.6), Inches(4.0),
        cd
    ).chart
    chart.has_legend = True
    chart.legend.position = -4131  # LEFT — corresponde ao modelo

    try:
        ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ser_el = list(chart._element.iter(f'{{{ns_c}}}ser'))[0]
        for i, cor in enumerate([AMBER, GREEN, NAVY]):
            dp = etree.SubElement(ser_el, f'{{{ns_c}}}dPt')
            idx = etree.SubElement(dp, f'{{{ns_c}}}idx')
            idx.set('val', str(i))
            spPr = etree.SubElement(dp, f'{{{ns_c}}}spPr')
            sf = etree.SubElement(
                etree.SubElement(spPr, f'{{{ns_a}}}solidFill'),
                f'{{{ns_a}}}srgbClr')
            sf.set('val', rgb_hex(cor))
    except Exception:
        pass

    show_datalabels(chart)
    set_chart_datalabel_size(chart, 13, True)

    vnc = [
        ("Pessoa Jurídica (PJ)", n_pj,  AMBER, "Prestadores/autônomos — não elegíveis CLT\nExcluídos de TotalPass, Wellhub e New Value"),
        ("Estagiário",           n_est, GREEN, "Excluídos do TotalPass e Wellhub\nNot excluídos do New Value"),
        ("Aprendiz",             n_apr, NAVY,  "Não atendem critérios de elegibilidade\nNão excluídos do New Value"),
    ]
    for i, (lbl, qtd, cor, desc) in enumerate(vnc):
        y = 1.25 + i * 1.32
        add_rect(s, 5.2, y, 4.3, 1.12, LGRAY)
        add_rect(s, 5.2, y, 0.07, 1.12, cor)
        pct = qtd / total_vinc * 100 if total_vinc else 0
        add_text(s, lbl, 5.35, y + 0.06, 2.4, 0.32, size=12, bold=True, color=NAVY)
        add_text(s, f"{qtd:,} exclusões ({pct:.1f}%)".replace(',', '.').replace('.', ',', 1),
                 5.35, y + 0.4, 2.4, 0.3, size=11, bold=True, color=cor)
        add_text(s, desc, 5.35, y + 0.72, 3.9, 0.36, size=9, color=GRAY)


# ── Slide 6: Empresas ─────────────────────────────────────────────────────────
def slide_empresas(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE

    add_text(s, "Exclusões por Empresa — Top 6", 0.5, 0.28, 9, 0.65,
             size=28, bold=True, color=NAVY, font="Cambria")
    add_text(s, "Empresas com maior volume de registros excluídos nos três processos",
             0.5, 0.88, 9, 0.3, size=12, color=GRAY)

    pe = data['por_empresa']
    nomes_curtos = {
        'STEFANINI CONSULTORIA E ASSESSORIA EM INFORMATICA S.A': 'Stefanini Consultoria',
        'W3HAUS COMUNICACAO INTERATIVA LTDA': 'W3Haus Comunicação',
        'GAUGE COMUNICACAO DIGITAL LTDA': 'Gauge Comunicação',
        'ORBITALL ATENDIMENTO LTDA': 'Orbitall Atendimento',
        'A11 TECNOLOGIA S.A.': 'A11 Tecnologia',
        'IHM': 'IHM',
    }
    items = [(nomes_curtos.get(k, k[:22]), v) for k, v in list(pe.items())[:6]]

    cd = CategoryChartData()
    cd.categories = [n for n, _ in items]
    cd.add_series('Exclusões', [v for _, v in items])

    cores_emp = [NAVY, BLUE, TEAL,
                 RGBColor(0x0E, 0xA5, 0xE9),
                 RGBColor(0x38, 0xBD, 0xF8),
                 RGBColor(0xBA, 0xE6, 0xFD)]

    chart = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.4), Inches(1.2), Inches(9.2), Inches(3.7),
        cd
    ).chart
    chart.has_legend = False

    try:
        ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ser_el = list(chart._element.iter(f'{{{ns_c}}}ser'))[0]
        for i, cor in enumerate(cores_emp):
            dp = etree.SubElement(ser_el, f'{{{ns_c}}}dPt')
            idx = etree.SubElement(dp, f'{{{ns_c}}}idx')
            idx.set('val', str(i))
            spPr = etree.SubElement(dp, f'{{{ns_c}}}spPr')
            sf = etree.SubElement(
                etree.SubElement(spPr, f'{{{ns_a}}}solidFill'),
                f'{{{ns_a}}}srgbClr')
            sf.set('val', rgb_hex(cor))
    except Exception:
        pass

    show_datalabels(chart)
    set_chart_datalabel_size(chart, 12, True)

    top2 = items[:2]
    soma_top2 = sum(v for _, v in top2)
    excl = data['total_excluidos']
    pct_top2 = soma_top2 / excl * 100 if excl else 0
    add_rect(s, 0.5, 5.08, 9.0, 0.38, LIGHT)
    add_text(s,
             f"{top2[0][0]} ({top2[0][1]:,}) + {top2[1][0]} ({top2[1][1]:,}) = {soma_top2:,} exclusões — {pct_top2:.1f}% do total".replace(',', '.').replace('.', ',', 2),
             0.65, 5.1, 8.6, 0.32, size=10, bold=True, color=NAVY)


# ── Slide 7: Cross Motivo × Processo ─────────────────────────────────────────
def slide_cross(prs, data):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)

    add_text(s, "Exclusões por Motivo e Processo", 0.5, 0.28, 9, 0.65,
             size=28, bold=True, color=NAVY, font="Cambria")
    add_text(s, "Quantidade de exclusões por categoria em cada um dos três processos",
             0.5, 0.88, 9, 0.3, size=12, color=GRAY)

    pm = data['por_motivo']
    pp = data['por_processo']

    labels = ["E-mail s/\ndomínio", "PJ", "E-mail\ninválido", "Estagiário", "Aprendiz", "Data\nfutura"]

    motivo_keys = [
        ["E-mail sem domínio Stefanini", "E-mail sem domínio IHM"],
        ["Pessoa Jurídica (PJ)"],
        ["E-mail inválido ou ausente"],
        ["Estagiário"],
        ["Aprendiz"],
        ["Data de admissão futura"],
    ]

    cross = data['cross']

    def get_cross(motivos, processo):
        total = 0
        for m in motivos:
            total += cross.get(processo, {}).get(m, 0)
        return total

    nv_vals = [get_cross(mk, 'New Value') for mk in motivo_keys]
    tp_vals = [get_cross(mk, 'TotalPass') for mk in motivo_keys]
    wh_vals = [get_cross(mk, 'Wellhub')   for mk in motivo_keys]

    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series('New Value', nv_vals)
    cd.add_series('TotalPass', tp_vals)
    cd.add_series('Wellhub',   wh_vals)

    chart = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.4), Inches(1.2), Inches(9.2), Inches(3.85),
        cd
    ).chart
    chart.has_legend = True
    chart.legend.position = -4160  # TOP

    cores_proc = [TEAL, NAVY, AMBER]
    try:
        ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        for i, (ser_el, cor) in enumerate(
                zip(chart._element.iter(f'{{{ns_c}}}ser'), cores_proc)):
            spPr = etree.SubElement(ser_el, f'{{{ns_c}}}spPr')
            sf = etree.SubElement(
                etree.SubElement(spPr, f'{{{ns_a}}}solidFill'),
                f'{{{ns_a}}}srgbClr')
            sf.set('val', rgb_hex(cor))
    except Exception:
        pass

    show_datalabels(chart)
    set_chart_datalabel_size(chart, 10, True)

    add_rect(s, 0.5, 5.08, 9.0, 0.38, RGBColor(0xEC, 0xFD, 0xF5))
    add_text(s, "New Value não exclui Estagiários nem Aprendizes — apenas Pessoa Jurídica (PJ).",
             0.65, 5.1, 8.6, 0.32, size=10, color=RGBColor(0x06, 0x5F, 0x46))


# ── Ponto de entrada ──────────────────────────────────────────────────────────
def generate(auditoria_path: str, output_path: str, total_base: int = 18218):
    stats = load_stats(auditoria_path)

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    slide_capa(prs, stats)
    slide_overview(prs, stats, total_base)
    slide_motivos(prs, stats)
    slide_email(prs, stats)
    slide_vinculo(prs, stats)
    slide_empresas(prs, stats)
    slide_cross(prs, stats)

    prs.save(output_path)
    return output_path